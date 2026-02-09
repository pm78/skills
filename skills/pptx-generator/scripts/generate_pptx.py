"""PPTX Generator - Creates presentations from JSON specs.

This module provides a PPTXGenerator class that creates PowerPoint presentations
from structured JSON configuration files. Supports multiple slide layouts including
title slides, bullet points, two-column layouts, image slides, and blank slides.
"""

from __future__ import annotations

import json
import math
import random
import re
import subprocess
import sys
import zipfile
import zlib
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Iterable, Optional, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from pptxgen import validate_config, validate_config_file


class PPTXGenerator:
    """Generate PPTX files from JSON specifications."""

    CHART_TYPES = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column-clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column-stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar-clustered": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar-stacked": XL_CHART_TYPE.BAR_STACKED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "donut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
    }

    LEGEND_POSITIONS = {
        "right": XL_LEGEND_POSITION.RIGHT,
        "left": XL_LEGEND_POSITION.LEFT,
        "top": XL_LEGEND_POSITION.TOP,
        "bottom": XL_LEGEND_POSITION.BOTTOM,
    }

    THEMES = {
        "dark": {
            "bg": RGBColor(15, 23, 42),  # slate-900
            "text_primary": RGBColor(241, 245, 249),  # slate-50
            "text_secondary": RGBColor(203, 213, 225),  # slate-300
            "accent": RGBColor(56, 189, 248),  # cyan-400
        },
        "light": {
            "bg": RGBColor(255, 255, 255),  # white
            "text_primary": RGBColor(15, 23, 42),  # slate-900
            "text_secondary": RGBColor(71, 85, 105),  # slate-600
            "accent": RGBColor(2, 132, 199),  # cyan-600
        },
    }

    SLIDE_WIDTH = Inches(10)
    SLIDE_HEIGHT = Inches(5.625)

    DIAGRAM_FALLBACK = {
        "fill": RGBColor(255, 255, 255),
        "line": RGBColor(2, 132, 199),  # cyan-600-ish
        "text_primary": RGBColor(15, 23, 42),
        "text_secondary": RGBColor(71, 85, 105),
        "arrow": RGBColor(2, 132, 199),
    }

    def __init__(
        self,
        config_path: str,
        theme: str = "dark",
        template_path: Optional[str] = None,
        assets_dir: Optional[str] = None,
        *,
        keep_template_slides: bool = False,
    ):
        self.config_path = Path(config_path).resolve()
        self.config_dir = self.config_path.parent
        self.assets_dir = Path(assets_dir).resolve() if assets_dir else None
        self.config, _ = validate_config_file(self.config_path)
        self._init_presentation(theme, template_path, keep_template_slides=keep_template_slides)

    @classmethod
    def from_dict(
        cls,
        config: Dict[str, Any],
        theme: str = "dark",
        template_path: Optional[str] = None,
        assets_dir: Optional[str] = None,
        *,
        keep_template_slides: bool = False,
        config_dir: Optional[Path] = None,
    ) -> "PPTXGenerator":
        instance = object.__new__(cls)
        instance.config, _ = validate_config(config)
        instance.config_path = None
        instance.config_dir = config_dir
        instance.assets_dir = Path(assets_dir).resolve() if assets_dir else None
        instance._init_presentation(theme, template_path, keep_template_slides=keep_template_slides)
        return instance

    def _init_presentation(self, theme: str, template_path: Optional[str], *, keep_template_slides: bool) -> None:
        self.theme_name = theme
        self.colors = None if theme == "template" else self.THEMES.get(theme, self.THEMES["dark"])

        self.template_path = Path(template_path) if template_path else None
        self.prs = Presentation(str(self.template_path)) if self.template_path else Presentation()
        self.template_theme_rgb = self._extract_template_theme_rgb() if self.template_path else {}

        if not self.template_path:
            self.prs.slide_width = self.SLIDE_WIDTH
            self.prs.slide_height = self.SLIDE_HEIGHT
        elif not keep_template_slides:
            self._remove_all_slides()

    def _remove_all_slides(self) -> None:
        # python-pptx has no public delete API; remove slide relationships directly.
        slide_id_list = self.prs.slides._sldIdLst  # type: ignore[attr-defined]
        for slide_id in list(slide_id_list):
            rel_id = slide_id.rId
            self.prs.part.drop_rel(rel_id)
            slide_id_list.remove(slide_id)

    def _extract_template_theme_rgb(self) -> Dict[MSO_THEME_COLOR, RGBColor]:
        if not self.template_path:
            return {}
        try:
            zf = zipfile.ZipFile(str(self.template_path))
        except Exception:
            return {}

        theme_xml_name = next(
            (name for name in zf.namelist() if name.startswith("ppt/theme/theme") and name.endswith(".xml")),
            None,
        )
        if not theme_xml_name:
            return {}

        try:
            raw = zf.read(theme_xml_name)
        except Exception:
            return {}

        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        try:
            root = ET.fromstring(raw)
        except Exception:
            return {}

        clr_scheme = root.find(".//a:clrScheme", ns)
        if clr_scheme is None:
            return {}

        key_map = {
            "bg1": MSO_THEME_COLOR.BACKGROUND_1,
            "bg2": MSO_THEME_COLOR.BACKGROUND_2,
            "tx1": MSO_THEME_COLOR.TEXT_1,
            "tx2": MSO_THEME_COLOR.TEXT_2,
            "dk1": MSO_THEME_COLOR.DARK_1,
            "lt1": MSO_THEME_COLOR.LIGHT_1,
            "dk2": MSO_THEME_COLOR.DARK_2,
            "lt2": MSO_THEME_COLOR.LIGHT_2,
            "accent1": MSO_THEME_COLOR.ACCENT_1,
            "accent2": MSO_THEME_COLOR.ACCENT_2,
            "accent3": MSO_THEME_COLOR.ACCENT_3,
            "accent4": MSO_THEME_COLOR.ACCENT_4,
            "accent5": MSO_THEME_COLOR.ACCENT_5,
            "accent6": MSO_THEME_COLOR.ACCENT_6,
            "hlink": MSO_THEME_COLOR.HYPERLINK,
            "folhlink": MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
        }

        def hex_to_rgb(value: str) -> Optional[RGBColor]:
            value = value.strip()
            if not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
                return None
            return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))

        result: Dict[MSO_THEME_COLOR, RGBColor] = {}
        for child in list(clr_scheme):
            tag = child.tag.split("}")[-1].lower()
            theme_color = key_map.get(tag)
            if not theme_color:
                continue

            srgb = child.find(".//a:srgbClr", ns)
            if srgb is not None:
                val = srgb.attrib.get("val")
                if val:
                    rgb = hex_to_rgb(val)
                    if rgb:
                        result[theme_color] = rgb
                        continue

            sys_clr = child.find(".//a:sysClr", ns)
            if sys_clr is not None:
                last = sys_clr.attrib.get("lastClr")
                if last:
                    rgb = hex_to_rgb(last)
                    if rgb:
                        result[theme_color] = rgb

        return result

    def _parse_theme_color(self, value: Any) -> Optional[MSO_THEME_COLOR]:
        if isinstance(value, MSO_THEME_COLOR):
            return value
        if value is None:
            return None
        raw = str(value).strip()
        if not raw:
            return None

        normalized = raw.upper().replace("-", "_").replace(" ", "_")
        normalized = re.sub(r"__+", "_", normalized)

        if normalized in MSO_THEME_COLOR.__members__:
            return MSO_THEME_COLOR[normalized]

        m = re.match(r"^(ACCENT|TEXT|BACKGROUND|DARK|LIGHT)_?(\d)$", normalized)
        if m:
            candidate = f"{m.group(1)}_{m.group(2)}"
            if candidate in MSO_THEME_COLOR.__members__:
                return MSO_THEME_COLOR[candidate]

        return None

    def _parse_rgb_color(self, value: Any) -> Optional[RGBColor]:
        if isinstance(value, RGBColor):
            return value
        if value is None:
            return None
        raw = str(value).strip()
        if not raw:
            return None

        if raw.startswith("#"):
            raw = raw[1:]
        if re.fullmatch(r"[0-9A-Fa-f]{6}", raw):
            return RGBColor(int(raw[0:2], 16), int(raw[2:4], 16), int(raw[4:6], 16))

        rgb_match = re.match(r"^rgb\\((\\d+)\\s*,\\s*(\\d+)\\s*,\\s*(\\d+)\\)$", raw, re.IGNORECASE)
        if rgb_match:
            r, g, b = (int(rgb_match.group(i)) for i in range(1, 4))
            return RGBColor(max(0, min(r, 255)), max(0, min(g, 255)), max(0, min(b, 255)))

        return None

    def _parse_color(self, value: Any) -> Optional[RGBColor | MSO_THEME_COLOR]:
        theme = self._parse_theme_color(value)
        if theme is not None:
            return theme
        return self._parse_rgb_color(value)

    def _theme_to_rgb(self, color: MSO_THEME_COLOR) -> RGBColor:
        if color in self.template_theme_rgb:
            return self.template_theme_rgb[color]
        # Fallback mapping when theme RGB extraction failed.
        fallback = {
            MSO_THEME_COLOR.ACCENT_1: self.DIAGRAM_FALLBACK["line"],
            MSO_THEME_COLOR.ACCENT_2: RGBColor(56, 189, 248),
            MSO_THEME_COLOR.ACCENT_3: RGBColor(2, 132, 199),
            MSO_THEME_COLOR.TEXT_1: self.DIAGRAM_FALLBACK["text_primary"],
            MSO_THEME_COLOR.TEXT_2: self.DIAGRAM_FALLBACK["text_secondary"],
            MSO_THEME_COLOR.BACKGROUND_1: self.DIAGRAM_FALLBACK["fill"],
        }
        return fallback.get(color, self.DIAGRAM_FALLBACK["line"])

    def _apply_color_to_fore(self, fore_color, color: RGBColor | MSO_THEME_COLOR) -> None:
        if isinstance(color, MSO_THEME_COLOR):
            fore_color.theme_color = color
        else:
            fore_color.rgb = color

    def _apply_color_to_font(self, font, color: RGBColor | MSO_THEME_COLOR) -> None:
        if isinstance(color, MSO_THEME_COLOR):
            font.color.theme_color = color
        else:
            font.color.rgb = color

    def _coerce_palette_section(self, palette: Any, section: str) -> Dict[str, Any]:
        if not isinstance(palette, dict):
            return {}
        if section in palette and isinstance(palette.get(section), dict):
            return dict(palette.get(section) or {})
        # Convenience: allow using diagram keys at the top-level.
        if section == "diagram" and any(k in palette for k in ("fill", "line", "text_primary", "text_secondary", "arrow")):
            return dict(palette)
        if section == "footer" and any(k in palette for k in ("text", "line")):
            return dict(palette)
        if section == "image" and any(k in palette for k in ("bg", "accent", "text")):
            return dict(palette)
        return {}

    def _slide_size_inches(self) -> Tuple[float, float]:
        w_in = float(int(self.prs.slide_width)) / float(Inches(1))
        h_in = float(int(self.prs.slide_height)) / float(Inches(1))
        return w_in, h_in

    def _format_human_date(self, dt: datetime, fmt: str) -> str:
        rendered = dt.strftime(fmt)
        # Common cleanup for day-of-month with leading zero in "%d".
        return rendered.replace(" 0", " ")

    def _resolve_deck_date_text(self, presentation_config: dict, professional_cfg: dict) -> str:
        date_cfg = professional_cfg.get("date")
        fallback = presentation_config.get("date")
        raw_value: Any = None
        date_fmt = "%B %d, %Y"

        if isinstance(date_cfg, str):
            raw_value = date_cfg
        elif isinstance(date_cfg, dict):
            raw_value = date_cfg.get("value")
            if isinstance(date_cfg.get("format"), str) and date_cfg.get("format").strip():
                date_fmt = str(date_cfg.get("format")).strip()

        if raw_value is None:
            if isinstance(fallback, str):
                raw_value = fallback
            elif isinstance(fallback, dict):
                raw_value = fallback.get("value")
                if isinstance(fallback.get("format"), str) and fallback.get("format").strip():
                    date_fmt = str(fallback.get("format")).strip()

        if raw_value is None or str(raw_value).strip() == "":
            return self._format_human_date(datetime.now(), date_fmt)

        raw = str(raw_value).strip()
        lowered = raw.lower()
        if lowered in {"today", "now"}:
            return self._format_human_date(datetime.now(), date_fmt)

        # Try to parse ISO-ish date strings; if parsing fails, keep raw text.
        try:
            parsed = datetime.fromisoformat(raw.replace("Z", "+00:00"))
            return self._format_human_date(parsed, date_fmt)
        except Exception:
            return raw

    def _resolve_professional_settings(self, presentation_config: dict) -> dict:
        raw = presentation_config.get("professional")
        if raw is False or raw is None:
            return {
                "enabled": False,
                "date_text": "",
                "date_on_title": False,
                "date_in_footer": False,
                "agenda_links": False,
                "force_agenda": False,
                "back_to_agenda": False,
                "back_to_agenda_label": "Back to Agenda",
            }

        if raw is True:
            raw = {}
        if not isinstance(raw, dict):
            raw = {}

        enabled = bool(raw.get("enabled", True))
        if not enabled:
            return {
                "enabled": False,
                "date_text": "",
                "date_on_title": False,
                "date_in_footer": False,
                "agenda_links": False,
                "force_agenda": False,
                "back_to_agenda": False,
                "back_to_agenda_label": "Back to Agenda",
            }

        date_cfg = raw.get("date")
        date_on_title = True
        date_in_footer = True
        if isinstance(date_cfg, dict):
            date_on_title = bool(date_cfg.get("on_title", True))
            date_in_footer = bool(date_cfg.get("in_footer", True))

        date_text = self._resolve_deck_date_text(presentation_config, raw)
        agenda_links = bool(raw.get("agenda_links", True))
        force_agenda = bool(raw.get("agenda", True))
        back_to_agenda = bool(raw.get("back_to_agenda", True))
        back_to_agenda_label = str(raw.get("back_to_agenda_label") or "Back to Agenda").strip() or "Back to Agenda"

        return {
            "enabled": True,
            "date_text": date_text,
            "date_on_title": date_on_title,
            "date_in_footer": date_in_footer,
            "agenda_links": agenda_links,
            "force_agenda": force_agenda,
            "back_to_agenda": back_to_agenda,
            "back_to_agenda_label": back_to_agenda_label,
        }

    def _add_cover_date(self, slide, *, date_text: str, slide_config: Optional[dict] = None) -> None:
        if not date_text:
            return

        w_in, h_in = self._slide_size_inches()
        palette = self._footer_palette(slide_config)
        box_w = 2.7
        box_h = 0.26
        box_x = max(0.5, w_in - box_w - 0.55)
        box_y = h_in - box_h - 0.28

        date_box = slide.shapes.add_textbox(Inches(box_x), Inches(box_y), Inches(box_w), Inches(box_h))
        tf = date_box.text_frame
        tf.clear()
        tf.word_wrap = False
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.02)
            tf.margin_right = Inches(0.02)
            tf.margin_top = Inches(0.01)
            tf.margin_bottom = Inches(0.01)
        except Exception:
            pass

        p = tf.paragraphs[0]
        p.text = date_text
        p.font.size = Pt(10 if (self.template_path and self.theme_name == "template") else 11)
        self._apply_color_to_font(p.font, palette["text"])
        p.alignment = PP_ALIGN.RIGHT

    def _add_back_to_agenda_link(
        self,
        slide,
        *,
        agenda_slide,
        label: str = "Back to Agenda",
        slide_config: Optional[dict] = None,
    ) -> None:
        if agenda_slide is None or not label:
            return

        w_in, _ = self._slide_size_inches()
        palette = self._footer_palette(slide_config)

        box_w = 1.7
        box_h = 0.22
        box_x = max(0.5, w_in - box_w - 0.55)
        box_y = 0.11 if (self.template_path and self.theme_name == "template") else 0.14

        link_box = slide.shapes.add_textbox(Inches(box_x), Inches(box_y), Inches(box_w), Inches(box_h))
        tf = link_box.text_frame
        tf.clear()
        tf.word_wrap = False
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Inches(0.01)
            tf.margin_right = Inches(0.01)
            tf.margin_top = Inches(0.01)
            tf.margin_bottom = Inches(0.01)
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(10)
        p.font.underline = True
        self._apply_color_to_font(p.font, palette["line"])
        p.alignment = PP_ALIGN.RIGHT

        try:
            link_box.click_action.target_slide = agenda_slide
        except Exception:
            pass

    def _should_footer(self, slide_config: dict, *, include_on_title: bool) -> bool:
        if slide_config.get("footer_exclude"):
            return False
        if slide_config.get("layout") == "title" and not include_on_title:
            return False
        return True

    def _footer_palette(self, slide_config: Optional[dict] = None) -> Dict[str, RGBColor | MSO_THEME_COLOR]:
        if self.colors:
            base: Dict[str, RGBColor | MSO_THEME_COLOR] = {
                "text": self.colors["text_secondary"],
                "line": self.colors["accent"],
            }
        else:
            base = {
                "text": MSO_THEME_COLOR.TEXT_2,
                "line": MSO_THEME_COLOR.ACCENT_1,
            }

        presentation_palette = getattr(self, "presentation_config", {}) or {}
        pres_section = self._coerce_palette_section(presentation_palette.get("palette"), "footer")
        slide_section = self._coerce_palette_section((slide_config or {}).get("palette"), "footer")
        merged = {**pres_section, **slide_section}
        for key in ("text", "line"):
            if key in merged:
                parsed = self._parse_color(merged[key])
                if parsed is not None:
                    base[key] = parsed

        return base

    def _add_footer(
        self,
        slide,
        *,
        doc_title: str,
        confidentiality: str,
        deck_date: Optional[str],
        page_no: Optional[int],
        total_pages: Optional[int],
        slide_config: Optional[dict] = None,
    ) -> None:
        w_in, h_in = self._slide_size_inches()
        palette = self._footer_palette(slide_config)

        # Template theme: keep the footer minimal and away from brand marks.
        # Worldline templates commonly include a bottom-left logo; so we place a
        # single right-aligned footer text box on the bottom-right.
        if self.template_path and self.theme_name == "template":
            parts = [t for t in [doc_title, confidentiality, deck_date] if t]
            if page_no is not None:
                if total_pages is not None and total_pages > 0:
                    parts.append(f"{page_no}/{total_pages}")
                else:
                    parts.append(str(page_no))
            footer_text = " • ".join(parts).strip()
            if not footer_text:
                return

            margin_right = 0.55
            margin_bottom = 0.18
            footer_h = 0.26
            footer_y = h_in - footer_h - margin_bottom
            box_w = 4.8
            box_x = max(0.5, w_in - margin_right - box_w)

            box = slide.shapes.add_textbox(Inches(box_x), Inches(footer_y), Inches(box_w), Inches(footer_h))
            tf = box.text_frame
            tf.clear()
            tf.word_wrap = False
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Inches(0.06)
                tf.margin_right = Inches(0.06)
                tf.margin_top = Inches(0.02)
                tf.margin_bottom = Inches(0.02)
            except Exception:
                pass
            p = tf.paragraphs[0]
            p.text = footer_text
            p.font.size = Pt(10)
            self._apply_color_to_font(p.font, palette["text"])
            p.alignment = PP_ALIGN.RIGHT
            return

        margin_x = 0.5
        footer_h = 0.28
        footer_y = h_in - footer_h - 0.06

        page_w = 1.0
        left_w = max(1.0, w_in - 2 * margin_x - page_w)
        right_x = w_in - margin_x - page_w

        # Divider line
        try:
            divider = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(margin_x),
                Inches(footer_y - 0.03),
                Inches(w_in - 2 * margin_x),
                Inches(0.01),
            )
            divider.line.width = Pt(0)
            divider.fill.solid()
            self._apply_color_to_fore(divider.fill.fore_color, palette["line"])
        except Exception:
            pass

        left_text = " | ".join([t for t in [doc_title, confidentiality, deck_date] if t])
        if left_text:
            left_box = slide.shapes.add_textbox(Inches(margin_x), Inches(footer_y), Inches(left_w), Inches(footer_h))
            tf = left_box.text_frame
            tf.clear()
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = left_text
            p.font.size = Pt(9)
            self._apply_color_to_font(p.font, palette["text"])
            p.alignment = PP_ALIGN.LEFT

        if page_no is not None:
            if total_pages is not None and total_pages > 0:
                page_text = f"{page_no}/{total_pages}"
            else:
                page_text = str(page_no)
            right_box = slide.shapes.add_textbox(Inches(right_x), Inches(footer_y), Inches(page_w), Inches(footer_h))
            tf = right_box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = page_text
            p.font.size = Pt(9)
            self._apply_color_to_font(p.font, palette["text"])
            p.alignment = PP_ALIGN.RIGHT

    def _insert_agenda_slide(self, slides: list[dict], *, with_links: bool, force_enabled: bool) -> list[dict]:
        if any(bool(s.get("is_agenda")) for s in slides):
            return slides

        cfg = getattr(self, "presentation_config", {}) or {}
        agenda_cfg = cfg.get("agenda")
        if agenda_cfg is False:
            if not force_enabled:
                return slides
            agenda_cfg = {}
        if agenda_cfg is True:
            agenda_cfg = {}
        if not isinstance(agenda_cfg, dict):
            agenda_cfg = {}

        enabled = bool(agenda_cfg.get("enabled", True))
        if force_enabled:
            enabled = True
        if not enabled:
            return slides
        if len(slides) <= 1:
            return slides

        entries: list[dict] = []
        for i, sc in enumerate(slides):
            if sc.get("layout") == "title":
                # Exclude title-style slides (cover/section breaks) from the agenda by default.
                continue
            if sc.get("agenda_exclude"):
                continue
            title = str(sc.get("agenda_title") or sc.get("title") or "").strip()
            if title:
                key = f"agenda-target-{i}"
                sc["_agenda_key"] = key
                entries.append({"text": title, "target_key": key})

        if not entries:
            return slides

        agenda_title = str(agenda_cfg.get("title") or "Agenda").strip()
        if with_links:
            agenda_slide = {
                "layout": "agenda",
                "title": agenda_title,
                "agenda_items": entries,
                "style": str(agenda_cfg.get("style") or "cards"),
                "is_agenda": True,
            }
            if "columns" in agenda_cfg:
                agenda_slide["columns"] = agenda_cfg.get("columns")
            if "template_layout" in agenda_cfg:
                agenda_slide["template_layout"] = agenda_cfg.get("template_layout")
            if "template_layout_index" in agenda_cfg:
                agenda_slide["template_layout_index"] = agenda_cfg.get("template_layout_index")
        else:
            titles = [str(item.get("text") or "").strip() for item in entries if str(item.get("text") or "").strip()]
            if len(titles) <= 10:
                agenda_slide = {"layout": "bullets", "title": agenda_title, "bullets": titles, "is_agenda": True}
            else:
                half = math.ceil(len(titles) / 2)
                left = "\n".join([f"• {t}" for t in titles[:half]])
                right = "\n".join([f"• {t}" for t in titles[half:]])
                agenda_slide = {
                    "layout": "two-column",
                    "title": agenda_title,
                    "left": left,
                    "right": right,
                    "is_agenda": True,
                }

        insert_at = 1 if slides and slides[0].get("layout") == "title" else 0
        return [*slides[:insert_at], agenda_slide, *slides[insert_at:]]

    def _resolve_agenda_links(self) -> None:
        pending = list(getattr(self, "_agenda_link_pending", []) or [])
        targets = dict(getattr(self, "_agenda_target_slides", {}) or {})
        for shape, target_key in pending:
            if not target_key:
                continue
            target_slide = targets.get(str(target_key))
            if target_slide is None:
                continue
            try:
                shape.click_action.target_slide = target_slide
            except Exception:
                continue

    def generate(self) -> Presentation:
        presentation_config = self.config.get("presentation", self.config)
        self.presentation_config = presentation_config
        slides = list(presentation_config.get("slides", []) or [])

        professional = self._resolve_professional_settings(presentation_config)
        deck_date = professional["date_text"] if professional.get("enabled") else ""

        # Agenda links are enabled by default in professional mode.
        agenda_cfg = presentation_config.get("agenda")
        agenda_links_requested = bool(professional.get("agenda_links", False))
        if isinstance(agenda_cfg, dict):
            if "links" in agenda_cfg:
                agenda_links_requested = bool(agenda_cfg.get("links"))
            elif "linked" in agenda_cfg:
                agenda_links_requested = bool(agenda_cfg.get("linked"))

        slides = self._insert_agenda_slide(
            slides,
            with_links=agenda_links_requested,
            force_enabled=bool(professional.get("force_agenda", False)),
        )
        presentation_config["slides"] = slides

        self._agenda_link_pending: list[tuple[Any, str]] = []
        self._agenda_target_slides: dict[str, Any] = {}

        layout_handlers = {
            "title": self._add_title_slide,
            "agenda": self._add_agenda_slide,
            "bullets": self._add_bullets_slide,
            "two-column": self._add_two_column_slide,
            "chart": self._add_chart_slide,
            "image": self._add_image_slide,
            "architecture": self._add_architecture_slide,
            "workflow": self._add_workflow_slide,
            "kpi-cards": self._add_kpi_cards_slide,
            "infographic": self._add_kpi_cards_slide,
            "blank": self._add_blank_slide,
        }

        footer_cfg = presentation_config.get("footer")
        if footer_cfg is False:
            footer_cfg = {"enabled": False}
        if not isinstance(footer_cfg, dict):
            footer_cfg = {}

        footer_enabled = bool(footer_cfg.get("enabled", True))
        doc_title = str(footer_cfg.get("title") or presentation_config.get("title") or "").strip()
        confidentiality = str(
            footer_cfg.get("confidentiality") or presentation_config.get("confidentiality") or "Confidential"
        ).strip()
        include_on_title = bool(footer_cfg.get("include_on_title", False))
        show_total = bool(footer_cfg.get("show_total", False))
        show_page_number = bool(footer_cfg.get("show_page_number", True))
        footer_date = deck_date if professional.get("date_in_footer") else ""

        total_pages = 0
        if footer_enabled and show_page_number:
            for sc in slides:
                if self._should_footer(sc, include_on_title=include_on_title):
                    total_pages += 1

        page_no = 0
        cover_date_added = False
        agenda_slide_ref = None
        for slide_config in slides:
            layout_type = slide_config.get("layout", "blank")
            handler = layout_handlers.get(layout_type, self._add_blank_slide)
            before = len(self.prs.slides)
            handler(slide_config)
            if len(self.prs.slides) <= before:
                continue
            slide = self.prs.slides[-1]

            # Agenda item targets are resolved after all slides are created.
            agenda_key = str(slide_config.get("_agenda_key") or "").strip()
            if agenda_key:
                self._agenda_target_slides[agenda_key] = slide

            if layout_type == "agenda" or bool(slide_config.get("is_agenda")):
                agenda_slide_ref = slide

            # Add date on the first rendered slide in professional mode.
            if professional.get("date_on_title") and deck_date and not cover_date_added:
                self._add_cover_date(slide, date_text=deck_date, slide_config=slide_config)
                cover_date_added = True

            if (
                professional.get("enabled")
                and professional.get("back_to_agenda")
                and agenda_slide_ref is not None
                and layout_type not in {"title", "agenda"}
                and not bool(slide_config.get("is_agenda"))
            ):
                self._add_back_to_agenda_link(
                    slide,
                    agenda_slide=agenda_slide_ref,
                    label=str(professional.get("back_to_agenda_label") or "Back to Agenda"),
                    slide_config=slide_config,
                )

            if footer_enabled and self._should_footer(slide_config, include_on_title=include_on_title):
                if show_page_number:
                    page_no += 1
                self._add_footer(
                    slide,
                    doc_title=doc_title,
                    confidentiality=confidentiality,
                    deck_date=footer_date,
                    page_no=page_no if show_page_number else None,
                    total_pages=total_pages if (show_page_number and show_total) else None,
                    slide_config=slide_config,
                )

        self._resolve_agenda_links()
        return self.prs

    def _get_blank_layout(self):
        blank_by_name = self._find_layout(["blank"])
        if blank_by_name is not None:
            return blank_by_name
        try:
            return self.prs.slide_layouts[6]
        except IndexError:
            return self.prs.slide_layouts[-1]

    def _set_slide_background(self, slide) -> None:
        if not self.colors:
            return
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.colors["bg"]

    def _add_title_textbox(self, slide, config: dict, y_pos: float = 0.4) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = config.get("title", "")
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        if self.colors:
            title_frame.paragraphs[0].font.color.rgb = self.colors["accent"]

    def _add_title_slide(self, config: dict) -> None:
        layout = self._resolve_layout(config, kind="title", name_candidates=["title slide", "title"])
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        if self.template_path and self.theme_name == "template":
            if self._set_title_placeholder(slide, config.get("title", "")):
                self._set_subtitle_placeholder(slide, config.get("subtitle", ""))
                return

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = config.get("title", "")
        title_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        if self.colors:
            title_frame.paragraphs[0].font.color.rgb = self.colors["accent"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        subtitle = config.get("subtitle")
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            subtitle_frame.paragraphs[0].font.size = Pt(28)
            if self.colors:
                subtitle_frame.paragraphs[0].font.color.rgb = self.colors["text_secondary"]
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_bullets_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="bullets",
            name_candidates=["title and content", "title and text", "content", "text"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        if self.template_path and self.theme_name == "template":
            title_set = self._set_title_placeholder(slide, config.get("title", ""))
            body = self._get_best_body_placeholder(slide)
            if body and hasattr(body, "text_frame"):
                self._fill_text_frame(body.text_frame, config.get("bullets", []), as_bullets=True)
                if config.get("title") and not title_set:
                    self._add_title_textbox(slide, config, y_pos=0.25)
                return

        self._add_title_textbox(slide, config)
        x, y, cx, cy = self._get_content_box(slide, has_title=True)
        bullet_box = slide.shapes.add_textbox(x, y, cx, cy)
        text_frame = bullet_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        bullets = config.get("bullets", [])
        for i, bullet in enumerate(bullets):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            normalized = bullet if bullet.lstrip().startswith("•") else f"• {bullet}"
            paragraph.text = normalized
            paragraph.font.size = Pt(20)
            if self.colors:
                paragraph.font.color.rgb = self.colors["text_primary"]
            paragraph.space_after = Pt(12)
            paragraph.level = 0

    def _add_two_column_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="two-column",
            name_candidates=["comparison", "two content", "two column", "two columns"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        force_manual_columns = bool(config.get("force_manual_columns"))
        title_set = False
        if self.template_path and self.theme_name == "template":
            title_set = self._set_title_placeholder(slide, config.get("title", ""))
            if not force_manual_columns and self._fill_two_column_placeholders(
                slide, config.get("left", ""), config.get("right", "")
            ):
                if config.get("title") and not title_set:
                    self._add_title_textbox(slide, config, y_pos=0.25)
                return

        if not (self.template_path and self.theme_name == "template" and title_set):
            self._add_title_textbox(slide, config)
        x, y, cx, cy = self._get_two_column_content_box(slide, has_title=True)
        gap = Inches(0.35)
        col_w = max(Inches(1.0), (int(cx) - int(gap)) // 2)
        left_box = slide.shapes.add_textbox(x, y, col_w, cy)
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        self._format_multiline_text(left_frame, config.get("left", ""))

        right_x = int(x) + int(col_w) + int(gap)
        right_box = slide.shapes.add_textbox(right_x, y, col_w, cy)
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        self._format_multiline_text(right_frame, config.get("right", ""))

    def _add_agenda_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="bullets",
            name_candidates=["agenda", "title and content", "title and text", "content", "text"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        title_set = False
        if self.template_path and self.theme_name == "template":
            title_set = self._set_title_placeholder(slide, config.get("title", ""))
        if not title_set:
            self._add_title_textbox(slide, config)

        # Remove placeholder scaffolding on agenda slides so edit mode doesn't show
        # "[Subtitle]" / "[Text]" prompts behind custom content.
        if self.template_path and self.theme_name == "template":
            removable_types = {
                PP_PLACEHOLDER.SUBTITLE,
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
                PP_PLACEHOLDER.VERTICAL_BODY,
                PP_PLACEHOLDER.VERTICAL_OBJECT,
            }
            title_shape = getattr(slide.shapes, "title", None)
            for ph in list(slide.placeholders):
                if ph is title_shape:
                    continue
                try:
                    ph_type = ph.placeholder_format.type
                except Exception:
                    continue
                if ph_type in removable_types:
                    try:
                        ph_el = ph._element
                        ph_el.getparent().remove(ph_el)
                    except Exception:
                        pass

        raw_items = config.get("agenda_items") if isinstance(config.get("agenda_items"), list) else []
        items: list[dict] = []
        for item in raw_items:
            if isinstance(item, dict):
                text = str(item.get("text") or "").strip()
                target_key = str(item.get("target_key") or "").strip()
            else:
                text = str(item).strip()
                target_key = ""
            if text:
                items.append({"text": text, "target_key": target_key})

        if not items:
            return
        if not hasattr(self, "_agenda_link_pending"):
            self._agenda_link_pending = []

        x, y, cx, cy = self._get_content_box(slide, has_title=True)
        style = str(config.get("style") or "cards").strip().lower()
        total = len(items)
        try:
            cols_raw = int(config.get("columns", 1 if total <= 8 else 2))
        except Exception:
            cols_raw = 1 if total <= 8 else 2
        cols = max(1, min(cols_raw, 3))
        rows = max(1, math.ceil(total / cols))

        if style == "cards":
            col_gap = int(Inches(0.28))
            item_gap = int(Inches(0.16))
            col_w = max(int(Inches(1.6)), (int(cx) - (cols - 1) * col_gap) // cols)
            row_h = max(int(Inches(0.42)), (int(cy) - (rows - 1) * item_gap) // rows)

            palette = self._diagram_palette(config)
            for idx, item in enumerate(items):
                col = idx // rows
                row = idx % rows
                if col >= cols:
                    col = cols - 1
                    row = rows - 1

                item_x = int(x) + col * (col_w + col_gap)
                item_y = int(y) + row * (row_h + item_gap)
                card = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    item_x,
                    item_y,
                    col_w,
                    row_h,
                )
                try:
                    card.fill.solid()
                    # Subtle surface tint keeps cards readable on both light and template themes.
                    fill_color = palette["fill"]
                    self._apply_color_to_fore(card.fill.fore_color, fill_color)
                    card.fill.transparency = 0.10
                except Exception:
                    pass
                try:
                    self._apply_color_to_fore(card.line.color, palette["line"])
                    card.line.width = Pt(1.1)
                except Exception:
                    pass

                tf = card.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                try:
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.margin_left = Inches(0.14)
                    tf.margin_right = Inches(0.14)
                    tf.margin_top = Inches(0.05)
                    tf.margin_bottom = Inches(0.05)
                except Exception:
                    pass

                item_label = str(item["text"]).strip()
                p = tf.paragraphs[0]
                p.text = f"{idx + 1:02d}  {item_label}"
                p.level = 0
                p.font.size = Pt(14 if (self.template_path and self.theme_name == "template") else 16)
                p.font.bold = True
                self._apply_color_to_font(p.font, palette["text_primary"])
                p.space_after = Pt(0)
                p.alignment = PP_ALIGN.LEFT

                target_key = item.get("target_key")
                if target_key:
                    self._agenda_link_pending.append((card, target_key))
        else:
            col_gap = int(Inches(0.35))
            item_gap = int(Inches(0.08))
            col_w = max(int(Inches(1.5)), (int(cx) - (cols - 1) * col_gap) // cols)
            row_h = max(int(Inches(0.34)), (int(cy) - (rows - 1) * item_gap) // rows)

            for idx, item in enumerate(items):
                col = idx // rows
                row = idx % rows
                if col >= cols:
                    col = cols - 1
                    row = rows - 1

                item_x = int(x) + col * (col_w + col_gap)
                item_y = int(y) + row * (row_h + item_gap)
                item_box = slide.shapes.add_textbox(item_x, item_y, col_w, row_h)
                tf = item_box.text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                try:
                    tf.vertical_anchor = MSO_ANCHOR.TOP
                    tf.margin_left = Inches(0.05)
                    tf.margin_right = Inches(0.05)
                    tf.margin_top = Inches(0.01)
                    tf.margin_bottom = Inches(0.01)
                except Exception:
                    pass

                p = tf.paragraphs[0]
                p.text = f"• {item['text']}"
                p.level = 0
                p.font.size = Pt(16 if (self.template_path and self.theme_name == "template") else 17)
                if self.colors:
                    p.font.color.rgb = self.colors["text_primary"]
                p.space_after = Pt(4)

                target_key = item.get("target_key")
                if target_key:
                    self._agenda_link_pending.append((item_box, target_key))

    def _get_two_column_content_box(self, slide, *, has_title: bool) -> Tuple[Any, Any, Any, Any]:
        # Templates can expose asymmetric body placeholders for "two-column" layouts.
        # Compute a stable union box across substantive body placeholders.
        if self.template_path and self.theme_name == "template":
            bodies = self._get_body_placeholders(slide)
            if bodies:
                substantial = [b for b in bodies if int(b.height) >= int(Inches(0.90))]
                if substantial:
                    left = min(int(b.left) for b in substantial)
                    top = min(int(b.top) for b in substantial)
                    right = max(int(b.left) + int(b.width) for b in substantial)
                    bottom = max(int(b.top) + int(b.height) for b in substantial)

                    # Keep content clear of title and footer/master elements.
                    title_shape = getattr(slide.shapes, "title", None)
                    if title_shape is not None:
                        top = max(top, int(title_shape.top) + int(title_shape.height) + int(Inches(0.10)))

                    max_bottom = int(self.prs.slide_height) - int(Inches(0.90))
                    bottom = min(bottom, max_bottom)
                    if bottom - top >= int(Inches(1.8)):
                        return left, top, max(int(Inches(1.6)), right - left), bottom - top

            # Conservative fallback box for template slides.
            y_offset = Inches(1.50) if has_title else Inches(0.85)
            reserve = Inches(0.90)
            return Inches(0.55), y_offset, Inches(8.9), max(Inches(1.8), self.prs.slide_height - y_offset - reserve)

        return self._get_content_box(slide, has_title=has_title)

    def _format_multiline_text(self, text_frame, text: str) -> None:
        text = self._normalize_newlines(text)
        lines = text.split("\n") if text else [""]

        for i, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            cleaned = line.strip()
            is_bullet = cleaned.startswith("•") or cleaned.startswith("- ") or cleaned.startswith("* ")
            if is_bullet:
                cleaned = cleaned.lstrip("•").lstrip("-").lstrip("*").lstrip()
                paragraph.text = f"• {cleaned}" if cleaned else "•"
            else:
                paragraph.text = cleaned

            base_size = 14 if (self.template_path and self.theme_name == "template") else 16
            paragraph.font.size = Pt(base_size)
            if i == 0 and not is_bullet and cleaned:
                paragraph.font.bold = True
                paragraph.font.size = Pt(base_size + 1)
            if self.colors:
                paragraph.font.color.rgb = self.colors["text_primary"]
            paragraph.space_after = Pt(6 if is_bullet else 8)

    def _diagram_palette(self, slide_config: Optional[dict] = None) -> Dict[str, RGBColor | MSO_THEME_COLOR]:
        if self.colors:
            if self.theme_name == "dark":
                base: Dict[str, RGBColor | MSO_THEME_COLOR] = {
                    "fill": RGBColor(30, 41, 59),  # slate-800
                    "line": self.colors["accent"],
                    "text_primary": self.colors["text_primary"],
                    "text_secondary": self.colors["text_secondary"],
                    "arrow": self.colors["accent"],
                }
            else:
                base = {
                    "fill": RGBColor(248, 250, 252),  # slate-50
                    "line": self.colors["accent"],
                    "text_primary": self.colors["text_primary"],
                    "text_secondary": self.colors["text_secondary"],
                    "arrow": self.colors["accent"],
                }
        else:
            # Template theme: use theme colors so fills/lines track the template palette.
            base = {
                "fill": MSO_THEME_COLOR.BACKGROUND_1,
                "line": MSO_THEME_COLOR.ACCENT_1,
                "text_primary": MSO_THEME_COLOR.TEXT_1,
                "text_secondary": MSO_THEME_COLOR.TEXT_2,
                "arrow": MSO_THEME_COLOR.ACCENT_1,
            }

        presentation_palette = getattr(self, "presentation_config", {}) or {}
        pres_section = self._coerce_palette_section(presentation_palette.get("palette"), "diagram")
        slide_section = self._coerce_palette_section((slide_config or {}).get("palette"), "diagram")
        merged = {**pres_section, **slide_section}
        for key in ("fill", "line", "text_primary", "text_secondary", "arrow"):
            if key in merged:
                parsed = self._parse_color(merged[key])
                if parsed is not None:
                    base[key] = parsed

        return base

    def _get_content_box(self, slide, *, has_title: bool) -> Tuple[Any, Any, Any, Any]:
        body = self._get_best_body_placeholder(slide)
        if body and hasattr(body, "left"):
            try:
                if hasattr(body, "text_frame"):
                    body.text_frame.clear()
            except Exception:
                pass
            reserve = int(Inches(0.75 if (self.template_path and self.theme_name == "template") else 0.45))
            height = max(int(Inches(0.6)), int(body.height) - reserve)
            return body.left, body.top, body.width, height

        y_offset = Inches(1.35) if has_title else Inches(0.7)
        base_h = Inches(3.9 if has_title else 4.5)
        reserve = Inches(0.75 if (self.template_path and self.theme_name == "template") else 0.45)
        return Inches(0.7), y_offset, Inches(8.6), max(Inches(0.6), base_h - reserve)

    def _add_workflow_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="workflow",
            name_candidates=["title and content", "title only", "title and text", "content"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        title = str(config.get("title") or "")
        if title:
            if self.template_path and self.theme_name == "template":
                if not self._set_title_placeholder(slide, title):
                    self._add_title_textbox(slide, {"title": title}, y_pos=0.25)
            else:
                self._add_title_textbox(slide, {"title": title}, y_pos=0.25)

        raw_steps = config.get("steps") or []
        steps: list[Dict[str, str]] = []
        for step in raw_steps:
            if isinstance(step, str):
                steps.append({"title": step})
                continue
            if isinstance(step, dict):
                step_title = str(step.get("title") or step.get("text") or step.get("label") or "").strip()
                if not step_title:
                    continue
                subtitle = str(step.get("subtitle") or step.get("note") or "").strip()
                item: Dict[str, str] = {"title": step_title}
                if subtitle:
                    item["subtitle"] = subtitle
                steps.append(item)

        if not steps:
            return

        orientation = str(config.get("orientation") or "horizontal").strip().lower()
        if orientation in {"v", "vert", "vertical"}:
            orientation = "vertical"
        else:
            orientation = "horizontal"

        show_numbers = bool(config.get("show_numbers", True))
        palette = self._diagram_palette(config)

        x0, y0, w0, h0 = self._get_content_box(slide, has_title=bool(title))
        x_in = float(int(x0)) / float(Inches(1))
        y_in = float(int(y0)) / float(Inches(1))
        w_in = float(int(w0)) / float(Inches(1))
        h_in = float(int(h0)) / float(Inches(1))

        if orientation == "horizontal":
            self._draw_workflow_horizontal(
                slide,
                steps=steps,
                x=x_in,
                y=y_in,
                w=w_in,
                h=h_in,
                palette=palette,
                show_numbers=show_numbers,
            )
        else:
            self._draw_workflow_vertical(
                slide,
                steps=steps,
                x=x_in,
                y=y_in,
                w=w_in,
                h=h_in,
                palette=palette,
                show_numbers=show_numbers,
            )

    def _add_architecture_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="architecture",
            name_candidates=["title and content", "title only", "title and text", "content"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        title = str(config.get("title") or "")
        if title:
            if self.template_path and self.theme_name == "template":
                if not self._set_title_placeholder(slide, title):
                    self._add_title_textbox(slide, {"title": title}, y_pos=0.25)
            else:
                self._add_title_textbox(slide, {"title": title}, y_pos=0.25)

        def normalize_nodes(value) -> list[Dict[str, str]]:
            nodes: list[Dict[str, str]] = []
            if not value:
                return nodes
            if not isinstance(value, list):
                value = [value]
            for item in value:
                if isinstance(item, str):
                    text = item.strip()
                    if text:
                        nodes.append({"title": text})
                    continue
                if isinstance(item, dict):
                    t = str(item.get("title") or item.get("text") or item.get("label") or "").strip()
                    if not t:
                        continue
                    sub = str(item.get("subtitle") or item.get("note") or "").strip()
                    node: Dict[str, str] = {"title": t}
                    if sub:
                        node["subtitle"] = sub
                    nodes.append(node)
            return nodes

        top_nodes = normalize_nodes(config.get("top_row") or config.get("top") or [])
        mid_nodes = normalize_nodes(config.get("middle_row") or config.get("middle") or [])
        bottom_node_list = normalize_nodes(config.get("bottom") or [])
        bottom_node = bottom_node_list[0] if bottom_node_list else {}

        if not top_nodes:
            top_nodes = [
                {"title": "React + Vite", "subtitle": "Frontend"},
                {"title": "FastAPI", "subtitle": "Backend APIs"},
                {"title": "LangGraph", "subtitle": "Orchestrator + Agents"},
            ]
        if not mid_nodes:
            mid_nodes = [
                {"title": "Vector Store", "subtitle": "FAISS • Pinecone • Chroma"},
                {"title": "Database", "subtitle": "SQLite • Postgres"},
                {"title": "Exports", "subtitle": "DOCX • PDF • HTML • PPTX"},
            ]
        if not bottom_node:
            bottom_node = {"title": "LLM Providers", "subtitle": "via LangChain (configurable)"}

        palette = self._diagram_palette(config)

        # Use a generous drawing area: full slide width minus margins, and reserve the
        # bottom for the Worldline logo and the footer.
        w_in, h_in = self._slide_size_inches()
        margin_x = 0.6
        top_y = 1.35 if title else 0.75
        bottom_reserved = 0.85 if (self.template_path and self.theme_name == "template") else 0.55
        x0 = margin_x
        y0 = top_y
        w0 = max(1.0, w_in - 2 * margin_x)
        h0 = max(1.0, h_in - top_y - bottom_reserved)

        # Row sizing (clamped to fit within h0).
        top_h = max(0.85, min(1.10, h0 * 0.34))
        mid_h = max(0.75, min(0.95, h0 * 0.26))
        gap_y = min(0.28, h0 * 0.08)
        remaining = h0 - top_h - gap_y - mid_h - gap_y
        bottom_h = max(0.75, min(1.20, remaining))
        if bottom_h + top_h + mid_h + 2 * gap_y > h0:
            # Tighten gaps as a fallback.
            gap_y = max(0.16, gap_y * 0.7)
            remaining = h0 - top_h - gap_y - mid_h - gap_y
            bottom_h = max(0.70, min(1.10, remaining))

        # Top row: 3 boxes + 2 arrows.
        arrow_w = 0.55
        gap_x = 0.16
        box_w = max(1.55, (w0 - 2 * arrow_w - 4 * gap_x) / 3)
        row_w = 3 * box_w + 2 * arrow_w + 4 * gap_x
        x_start = x0 + max(0.0, (w0 - row_w) / 2)
        y_top = y0

        xs = [x_start, x_start + box_w + gap_x + arrow_w + gap_x, x_start + 2 * (box_w + gap_x + arrow_w + gap_x)]
        for i in range(3):
            node = top_nodes[i] if i < len(top_nodes) else {}
            self._draw_workflow_step_box(
                slide,
                x=xs[i],
                y=y_top,
                w=box_w,
                h=top_h,
                title=str(node.get("title") or ""),
                subtitle=str(node.get("subtitle") or ""),
                step_number=None,
                palette=palette,
            )

        arrow_h = min(0.42, top_h * 0.35)
        y_arrow = y_top + (top_h - arrow_h) / 2
        for i in range(2):
            x_arrow = xs[i] + box_w + gap_x
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x_arrow),
                Inches(y_arrow),
                Inches(arrow_w),
                Inches(arrow_h),
            )
            self._style_arrow(arrow, palette)

        # Middle row: 3 boxes, slightly smaller.
        y_mid = y_top + top_h + gap_y
        mid_gap_x = 0.30
        mid_box_w = max(1.45, (w0 - 2 * mid_gap_x) / 3)
        mid_row_w = 3 * mid_box_w + 2 * mid_gap_x
        mid_x_start = x0 + max(0.0, (w0 - mid_row_w) / 2)
        mid_xs = [mid_x_start, mid_x_start + mid_box_w + mid_gap_x, mid_x_start + 2 * (mid_box_w + mid_gap_x)]
        for i in range(3):
            node = mid_nodes[i] if i < len(mid_nodes) else {}
            self._draw_workflow_step_box(
                slide,
                x=mid_xs[i],
                y=y_mid,
                w=mid_box_w,
                h=mid_h,
                title=str(node.get("title") or ""),
                subtitle=str(node.get("subtitle") or ""),
                step_number=None,
                palette=palette,
            )

        # Bottom row: wide box.
        y_bottom = y_mid + mid_h + gap_y
        bottom_w = w0
        bottom_x = x0
        self._draw_workflow_step_box(
            slide,
            x=bottom_x,
            y=y_bottom,
            w=bottom_w,
            h=bottom_h,
            title=str(bottom_node.get("title") or ""),
            subtitle=str(bottom_node.get("subtitle") or ""),
            step_number=None,
            palette=palette,
        )

    def _draw_workflow_horizontal(
        self,
        slide,
        *,
        steps: list[Dict[str, str]],
        x: float,
        y: float,
        w: float,
        h: float,
        palette: Dict[str, RGBColor | MSO_THEME_COLOR],
        show_numbers: bool,
    ) -> None:
        n = len(steps)
        if n == 0:
            return

        min_box_w = 1.35
        arrow_w = 0.45
        gap = 0.10
        max_per_row = min(n, 6)

        def box_width(per_row: int, *, arrow_w_in: float, gap_in: float) -> float:
            if per_row <= 1:
                return w
            return (w - (per_row - 1) * (arrow_w_in + 2 * gap_in)) / per_row

        chosen = 1
        for candidate in range(max_per_row, 1, -1):
            bw = box_width(candidate, arrow_w_in=arrow_w, gap_in=gap)
            if bw >= min_box_w:
                chosen = candidate
                break
        if chosen == 1:
            # Try tightening arrows/gaps before falling back to 1 per row.
            for candidate in range(max_per_row, 1, -1):
                bw = box_width(candidate, arrow_w_in=0.34, gap_in=0.06)
                if bw >= min_box_w:
                    chosen = candidate
                    arrow_w = 0.34
                    gap = 0.06
                    break

        per_row = max(1, chosen)
        rows = max(1, math.ceil(n / per_row))
        row_gap = 0.35
        box_h = min(1.35, (h - (rows - 1) * row_gap) / rows)
        box_h = max(0.9, box_h)

        idx = 0
        for row in range(rows):
            remaining = n - idx
            count = min(per_row, remaining)
            if count <= 0:
                break

            direction = 1 if row % 2 == 0 else -1  # serpentine
            bw = box_width(count, arrow_w_in=arrow_w, gap_in=gap)
            step_stride = bw + (arrow_w + 2 * gap)

            y_box = y + row * (box_h + row_gap)
            xs: list[float] = []
            if direction == 1:
                xs = [x + j * step_stride for j in range(count)]
            else:
                xs = [x + w - bw - j * step_stride for j in range(count)]

            # Boxes
            for j in range(count):
                step = steps[idx + j]
                self._draw_workflow_step_box(
                    slide,
                    x=xs[j],
                    y=y_box,
                    w=bw,
                    h=box_h,
                    title=step.get("title", ""),
                    subtitle=step.get("subtitle", ""),
                    step_number=(idx + j + 1) if show_numbers else None,
                    palette=palette,
                )

            # In-row arrows
            arrow_shape = MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW if direction == 1 else MSO_AUTO_SHAPE_TYPE.LEFT_ARROW
            arrow_h = min(0.45, box_h * 0.35)
            y_arrow = y_box + (box_h - arrow_h) / 2
            for j in range(count - 1):
                if direction == 1:
                    x_arrow = xs[j] + bw + gap
                else:
                    x_arrow = xs[j + 1] + bw + gap
                arrow = slide.shapes.add_shape(
                    arrow_shape,
                    Inches(x_arrow),
                    Inches(y_arrow),
                    Inches(arrow_w),
                    Inches(arrow_h),
                )
                self._style_arrow(arrow, palette)

            # Between-row arrow (down)
            if row < rows - 1 and (idx + count) < n:
                end_x = xs[-1] if direction == 1 else xs[0]
                down_w = min(0.55, bw * 0.35)
                down_h = min(0.40, row_gap * 0.9)
                x_down = end_x + (bw - down_w) / 2
                y_down = y_box + box_h + (row_gap - down_h) / 2
                down = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                    Inches(x_down),
                    Inches(y_down),
                    Inches(down_w),
                    Inches(down_h),
                )
                self._style_arrow(down, palette)

            idx += count

    def _draw_workflow_vertical(
        self,
        slide,
        *,
        steps: list[Dict[str, str]],
        x: float,
        y: float,
        w: float,
        h: float,
        palette: Dict[str, RGBColor | MSO_THEME_COLOR],
        show_numbers: bool,
    ) -> None:
        n = len(steps)
        if n == 0:
            return

        min_box_h = 0.75
        arrow_h = 0.28
        gap = 0.08
        max_per_col = min(n, 8)

        def box_height(per_col: int, *, arrow_h_in: float, gap_in: float) -> float:
            if per_col <= 1:
                return h
            return (h - (per_col - 1) * (arrow_h_in + 2 * gap_in)) / per_col

        chosen = 1
        for candidate in range(max_per_col, 1, -1):
            bh = box_height(candidate, arrow_h_in=arrow_h, gap_in=gap)
            if bh >= min_box_h:
                chosen = candidate
                break

        per_col = max(1, chosen)
        cols = max(1, math.ceil(n / per_col))
        col_gap = 0.35
        box_w = min(2.8, (w - (cols - 1) * col_gap) / cols)
        box_w = max(1.8, box_w)

        idx = 0
        for col in range(cols):
            remaining = n - idx
            count = min(per_col, remaining)
            if count <= 0:
                break

            direction = 1 if col % 2 == 0 else -1  # serpentine
            bh = box_height(count, arrow_h_in=arrow_h, gap_in=gap)
            step_stride = bh + (arrow_h + 2 * gap)

            x_box = x + col * (box_w + col_gap)
            ys: list[float] = []
            if direction == 1:
                ys = [y + j * step_stride for j in range(count)]
            else:
                ys = [y + h - bh - j * step_stride for j in range(count)]

            for j in range(count):
                step = steps[idx + j]
                self._draw_workflow_step_box(
                    slide,
                    x=x_box,
                    y=ys[j],
                    w=box_w,
                    h=bh,
                    title=step.get("title", ""),
                    subtitle=step.get("subtitle", ""),
                    step_number=(idx + j + 1) if show_numbers else None,
                    palette=palette,
                )

            arrow_shape = MSO_AUTO_SHAPE_TYPE.DOWN_ARROW if direction == 1 else MSO_AUTO_SHAPE_TYPE.UP_ARROW
            aw = min(0.55, box_w * 0.30)
            for j in range(count - 1):
                if direction == 1:
                    y_arrow = ys[j] + bh + gap
                else:
                    y_arrow = ys[j + 1] + bh + gap
                x_arrow = x_box + (box_w - aw) / 2
                arr = slide.shapes.add_shape(
                    arrow_shape,
                    Inches(x_arrow),
                    Inches(y_arrow),
                    Inches(aw),
                    Inches(arrow_h),
                )
                self._style_arrow(arr, palette)

            if col < cols - 1 and (idx + count) < n:
                end_y = ys[-1] if direction == 1 else ys[0]
                rw = min(0.55, col_gap * 0.9)
                rh = min(0.45, bh * 0.35)
                x_right = x_box + box_w + (col_gap - rw) / 2
                y_right = end_y + (bh - rh) / 2
                right = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                    Inches(x_right),
                    Inches(y_right),
                    Inches(rw),
                    Inches(rh),
                )
                self._style_arrow(right, palette)

            idx += count

    def _style_arrow(self, shape, palette: Dict[str, RGBColor | MSO_THEME_COLOR]) -> None:
        try:
            shape.fill.solid()
            self._apply_color_to_fore(shape.fill.fore_color, palette["arrow"])
        except Exception:
            pass
        try:
            self._apply_color_to_fore(shape.line.color, palette["arrow"])
            shape.line.width = Pt(1)
        except Exception:
            pass

    def _draw_workflow_step_box(
        self,
        slide,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        title: str,
        subtitle: str,
        step_number: Optional[int],
        palette: Dict[str, RGBColor | MSO_THEME_COLOR],
    ) -> None:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        try:
            box.fill.solid()
            self._apply_color_to_fore(box.fill.fore_color, palette["fill"])
        except Exception:
            pass
        try:
            self._apply_color_to_fore(box.line.color, palette["line"])
            box.line.width = Pt(1.25)
        except Exception:
            pass

        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.12)
            tf.margin_right = Inches(0.12)
            tf.margin_top = Inches(0.10)
            tf.margin_bottom = Inches(0.10)
        except Exception:
            pass

        header = title.strip()
        if step_number is not None:
            header = f"{step_number}. {header}" if header else f"{step_number}."

        p1 = tf.paragraphs[0]
        p1.text = header
        p1.alignment = PP_ALIGN.CENTER
        p1.font.bold = True
        p1.font.size = Pt(16)
        self._apply_color_to_font(p1.font, palette["text_primary"])

        if subtitle.strip():
            p2 = tf.add_paragraph()
            p2.text = subtitle.strip()
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = Pt(12)
            self._apply_color_to_font(p2.font, palette["text_secondary"])

    def _add_kpi_cards_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="kpi-cards",
            name_candidates=["title and content", "title only", "title and text", "content"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        title = str(config.get("title") or "")
        if title:
            if self.template_path and self.theme_name == "template":
                if not self._set_title_placeholder(slide, title):
                    self._add_title_textbox(slide, {"title": title}, y_pos=0.25)
            else:
                self._add_title_textbox(slide, {"title": title}, y_pos=0.25)

        raw_cards = config.get("cards") or config.get("items") or []
        cards: list[Dict[str, str]] = []
        for item in raw_cards:
            if isinstance(item, str):
                cards.append({"label": item})
                continue
            if isinstance(item, dict):
                label = str(item.get("label") or item.get("title") or "").strip()
                value = str(item.get("value") or item.get("metric") or item.get("number") or "").strip()
                note = str(item.get("note") or item.get("subtitle") or "").strip()
                icon_path = str(item.get("icon_path") or item.get("image_path") or "").strip()
                card: Dict[str, str] = {}
                if label:
                    card["label"] = label
                if value:
                    card["value"] = value
                if note:
                    card["note"] = note
                if icon_path:
                    card["icon_path"] = icon_path
                if card:
                    cards.append(card)

        if not cards:
            return

        x0, y0, w0, h0 = self._get_content_box(slide, has_title=bool(title))
        x_in = float(int(x0)) / float(Inches(1))
        y_in = float(int(y0)) / float(Inches(1))
        w_in = float(int(w0)) / float(Inches(1))
        h_in = float(int(h0)) / float(Inches(1))

        columns = config.get("columns")
        if columns is None:
            columns = 2 if len(cards) <= 6 else 3
        try:
            cols = max(1, min(int(columns), len(cards)))
        except Exception:
            cols = 2
        rows = max(1, math.ceil(len(cards) / cols))

        gap_x = 0.28
        gap_y = 0.24
        card_w = (w_in - (cols - 1) * gap_x) / cols
        card_h = (h_in - (rows - 1) * gap_y) / rows

        palette = self._diagram_palette(config)

        for i, card in enumerate(cards):
            r = i // cols
            c = i % cols
            x_card = x_in + c * (card_w + gap_x)
            y_card = y_in + r * (card_h + gap_y)
            self._draw_kpi_card(
                slide,
                x=x_card,
                y=y_card,
                w=card_w,
                h=card_h,
                label=card.get("label", ""),
                value=card.get("value", ""),
                note=card.get("note", ""),
                icon_path=card.get("icon_path"),
                palette=palette,
            )

    def _draw_kpi_card(
        self,
        slide,
        *,
        x: float,
        y: float,
        w: float,
        h: float,
        label: str,
        value: str,
        note: str,
        icon_path: Optional[str],
        palette: Dict[str, RGBColor | MSO_THEME_COLOR],
    ) -> None:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        try:
            card.fill.solid()
            self._apply_color_to_fore(card.fill.fore_color, palette["fill"])
        except Exception:
            pass
        try:
            self._apply_color_to_fore(card.line.color, palette["line"])
            card.line.width = Pt(1.25)
        except Exception:
            pass

        # Optional icon (top-right)
        if icon_path:
            img = self._resolve_fs_path(icon_path)
            if img and img.exists():
                size = min(0.55, h * 0.32)
                margin = 0.12
                slide.shapes.add_picture(
                    str(img),
                    Inches(x + w - size - margin),
                    Inches(y + margin),
                    width=Inches(size),
                    height=Inches(size),
                )

        tf = card.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.18)
            tf.margin_right = Inches(0.18)
            tf.margin_top = Inches(0.14)
            tf.margin_bottom = Inches(0.10)
        except Exception:
            pass

        # Value (big)
        if value:
            p0 = tf.paragraphs[0]
            p0.text = value
            value_len = len(value.strip())
            if value_len >= 12:
                value_size = 18
            elif value_len >= 9:
                value_size = 22
            else:
                value_size = 28
            if self.template_path and self.theme_name == "template":
                value_size = max(16, value_size - 3)
            p0.font.size = Pt(value_size)
            p0.font.bold = True
            self._apply_color_to_font(p0.font, palette["text_primary"])
        else:
            tf.paragraphs[0].text = ""

        # Label (medium)
        if label:
            p1 = tf.add_paragraph()
            p1.text = label
            p1.font.size = Pt(12 if (self.template_path and self.theme_name == "template") else 14)
            p1.font.bold = True
            self._apply_color_to_font(p1.font, palette["text_secondary"])

        # Note (small)
        if note:
            p2 = tf.add_paragraph()
            p2.text = note
            p2.font.size = Pt(10 if (self.template_path and self.theme_name == "template") else 12)
            self._apply_color_to_font(p2.font, palette["text_secondary"])

    def _add_image_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="image",
            name_candidates=["picture with caption", "picture", "image"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        if config.get("title"):
            if not (self.template_path and self.theme_name == "template"):
                self._add_title_textbox(slide, config, y_pos=0.2)
            else:
                if not self._set_title_placeholder(slide, config.get("title", "")):
                    self._add_title_textbox(slide, config, y_pos=0.2)

        if self.template_path and self.theme_name == "template":
            usage_subtitle = str(config.get("usage_subtitle") or "").strip()
            usage_points = config.get("usage_points") if isinstance(config.get("usage_points"), list) else []
            usage_lines = [str(item).strip() for item in usage_points if str(item).strip()]
            force_manual_usage_text = bool(config.get("force_manual_usage_text"))
            usage_as_bullets = bool(config.get("usage_as_bullets", True))

            body_types = {
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
                PP_PLACEHOLDER.VERTICAL_BODY,
                PP_PLACEHOLDER.VERTICAL_OBJECT,
            }
            body_placeholders = []
            for placeholder in slide.placeholders:
                if not hasattr(placeholder, "text_frame"):
                    continue
                try:
                    ph_type = placeholder.placeholder_format.type
                    if ph_type not in body_types:
                        continue
                except Exception:
                    pass
                body_placeholders.append(placeholder)

            body_placeholders = sorted(body_placeholders, key=lambda ph: int(ph.width) * int(ph.height), reverse=True)

            subtitle_placeholder = None
            main_body_placeholder = body_placeholders[0] if body_placeholders else None
            if len(body_placeholders) >= 2:
                subtitle_placeholder = min(body_placeholders, key=lambda ph: int(ph.top))
                if subtitle_placeholder is main_body_placeholder:
                    subtitle_placeholder = body_placeholders[1]

            if force_manual_usage_text and (usage_subtitle or usage_lines):
                for placeholder in body_placeholders:
                    try:
                        placeholder.text_frame.clear()
                    except Exception:
                        pass

                slide_mid = int(self.prs.slide_width) / 2
                left_ph = [ph for ph in body_placeholders if (int(ph.left) + int(ph.width) / 2) < slide_mid]
                if left_ph:
                    left = min(int(ph.left) for ph in left_ph)
                    top = min(int(ph.top) for ph in left_ph)
                    right = max(int(ph.left) + int(ph.width) for ph in left_ph)
                    bottom = max(int(ph.top) + int(ph.height) for ph in left_ph)
                else:
                    left = int(Inches(0.55))
                    top = int(Inches(1.45))
                    right = int(Inches(4.45))
                    bottom = int(self.prs.slide_height - Inches(0.95))

                title_shape = getattr(slide.shapes, "title", None)
                if title_shape is not None:
                    title_bottom = int(title_shape.top) + int(title_shape.height) + int(Inches(0.08))
                    top = max(top, title_bottom)

                bottom = min(bottom, int(self.prs.slide_height) - int(Inches(0.95)))
                if bottom - top < int(Inches(1.6)):
                    top = int(Inches(1.45))
                    bottom = int(self.prs.slide_height) - int(Inches(0.95))

                usage_title_h = int(Inches(0.42)) if usage_subtitle else 0
                if usage_subtitle:
                    subtitle_box = slide.shapes.add_textbox(left, top, right - left, usage_title_h)
                    subtitle_tf = subtitle_box.text_frame
                    subtitle_tf.clear()
                    subtitle_tf.word_wrap = True
                    subtitle_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    p = subtitle_tf.paragraphs[0]
                    p.text = usage_subtitle
                    p.font.bold = True
                    p.font.size = Pt(22 if (self.template_path and self.theme_name == "template") else 20)
                    if self.colors:
                        p.font.color.rgb = self.colors["text_primary"]
                    top = top + usage_title_h + int(Inches(0.05))

                if usage_lines:
                    usage_box = slide.shapes.add_textbox(left, top, right - left, max(int(Inches(1.0)), bottom - top))
                    usage_tf = usage_box.text_frame
                    usage_tf.word_wrap = True
                    usage_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    self._fill_text_frame(usage_tf, usage_lines, as_bullets=usage_as_bullets)
            else:
                if usage_subtitle:
                    if not self._set_subtitle_placeholder(slide, usage_subtitle):
                        if subtitle_placeholder is not None:
                            self._fill_text_frame(subtitle_placeholder.text_frame, [usage_subtitle], as_bullets=False)

                if usage_lines and main_body_placeholder is not None:
                    self._fill_text_frame(main_body_placeholder.text_frame, usage_lines, as_bullets=usage_as_bullets)

                for placeholder in body_placeholders:
                    if placeholder is main_body_placeholder or placeholder is subtitle_placeholder:
                        continue
                    try:
                        placeholder.text_frame.clear()
                    except Exception:
                        pass

        img_path = self._resolve_fs_path(config.get("image_path"))
        if not (img_path and img_path.exists()):
            img_path = self._generate_image_asset(config)
        elif img_path:
            img_path = self._maybe_autocrop_image(img_path, config)

        if img_path and img_path.exists():
            if self.template_path and self.theme_name == "template":
                picture_ph = None
                try:
                    candidates = [ph for ph in slide.placeholders if hasattr(ph, "insert_picture")]
                    if candidates:
                        picture_ph = max(candidates, key=lambda ph: int(ph.width) * int(ph.height))
                except Exception:
                    picture_ph = None

                if picture_ph:
                    ph_left = int(picture_ph.left)
                    ph_top = int(picture_ph.top)
                    ph_width = int(picture_ph.width)
                    ph_height = int(picture_ph.height)
                    pic = picture_ph.insert_picture(str(img_path))
                    try:
                        pic.crop_left = 0
                        pic.crop_right = 0
                        pic.crop_top = 0
                        pic.crop_bottom = 0
                    except Exception:
                        pass
                    try:
                        # Keep template picture clear of footer text/master marks.
                        max_bottom = int(self.prs.slide_height) - int(Inches(0.88))
                        overflow = int(pic.top) + int(pic.height) - max_bottom
                        if overflow > 0:
                            pic.height = max(int(Inches(1.0)), int(pic.height) - overflow)
                    except Exception:
                        pass

                    # Some templates produce invalid picture geometry (e.g., zero width) after insert_picture.
                    # Fallback: manually place a contained image inside the picture placeholder bounds.
                    invalid_geometry = False
                    try:
                        invalid_geometry = int(pic.width) < int(Inches(0.5)) or int(pic.height) < int(Inches(0.5))
                    except Exception:
                        invalid_geometry = True
                    if invalid_geometry:
                        safe_top = max(0, ph_top)
                        safe_height = ph_height - max(0, safe_top - ph_top)
                        if safe_height > int(Inches(0.8)) and ph_width > int(Inches(1.0)):
                            safe_bottom = min(int(self.prs.slide_height) - int(Inches(0.88)), safe_top + safe_height)
                            safe_height = max(int(Inches(0.8)), safe_bottom - safe_top)
                            safe_left = ph_left
                            self._add_picture_contain(
                                slide,
                                img_path=img_path,
                                x=safe_left,
                                y=safe_top,
                                cx=ph_width,
                                cy=safe_height,
                            )

                    # Optional override for templates where the picture placeholder is too small.
                    if bool(config.get("expand_image_to_content_box")):
                        x, y, cx, cy = self._get_content_box(slide, has_title=bool(config.get("title")))
                        left, top, w, h = self._compute_contain_geometry(img_path=img_path, x=x, y=y, cx=cx, cy=cy)
                        try:
                            pic.left = left
                            pic.top = top
                            pic.width = w
                            pic.height = h
                        except Exception:
                            self._add_picture_contain(slide, img_path=img_path, x=x, y=y, cx=cx, cy=cy)
                else:
                    x, y, cx, cy = self._get_content_box(slide, has_title=bool(config.get("title")))
                    self._add_picture_contain(slide, img_path=img_path, x=x, y=y, cx=cx, cy=cy)
            else:
                x, y, cx, cy = self._get_content_box(slide, has_title=bool(config.get("title")))
                self._add_picture_contain(slide, img_path=img_path, x=x, y=y, cx=cx, cy=cy)

        caption = config.get("caption")
        if caption:
            w_in, h_in = self._slide_size_inches()
            # Keep caption above footer/logo area for templates.
            if self.template_path and self.theme_name == "template":
                caption_y = h_in - 1.20
                cap_h = 0.28
            else:
                caption_y = h_in - 0.65
                cap_h = 0.35
            caption_box = slide.shapes.add_textbox(Inches(0.5), Inches(caption_y), Inches(w_in - 1.0), Inches(cap_h))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            caption_frame.paragraphs[0].font.size = Pt(12)
            caption_frame.paragraphs[0].font.italic = True
            if self.colors:
                caption_frame.paragraphs[0].font.color.rgb = self.colors["text_secondary"]
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_blank_slide(self, config: dict) -> None:
        slide = self.prs.slides.add_slide(self._resolve_layout(config, kind="blank", name_candidates=["blank"]))
        self._set_slide_background(slide)

        if config.get("title"):
            if self.template_path and self.theme_name == "template":
                if not self._set_title_placeholder(slide, config.get("title", "")):
                    self._add_title_textbox(slide, config, y_pos=0.25)
            else:
                self._add_title_textbox(slide, config)

    def save(self, output_path: str) -> Path:
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output))
        print(f"✅ PPTX saved to {output}")
        return output

    def _find_layout(self, name_candidates: Iterable[str]):
        candidates = [c.strip().lower() for c in name_candidates if c.strip()]
        if not candidates:
            return None

        def normalize(name: str) -> str:
            return re.sub(r"[^a-z0-9]+", " ", (name or "").strip().lower()).strip()

        for slide_layout in self.prs.slide_layouts:
            name = getattr(slide_layout, "name", "") or ""
            normalized = name.strip().lower()
            if any(normalized == c for c in candidates):
                return slide_layout

        best_layout = None
        best_score = -1
        for slide_layout in self.prs.slide_layouts:
            name = getattr(slide_layout, "name", "") or ""
            normalized = normalize(name)

            score = 0
            for candidate in candidates:
                c = normalize(candidate)
                if not c:
                    continue
                if normalized == c:
                    score = max(score, 200 + len(c))
                    continue
                if re.search(rf"\b{re.escape(c)}\b", normalized):
                    score = max(score, 120 + len(c))
                    continue
                if c in normalized:
                    score = max(score, 80 + len(c))

            if score > best_score:
                best_layout = slide_layout
                best_score = score

        if best_score > 0:
            return best_layout

        return None

    def _resolve_layout(self, slide_config: dict, *, kind: str, name_candidates: Iterable[str]):
        if not self.template_path:
            return self._get_blank_layout()

        layout_index = slide_config.get("template_layout_index")
        if layout_index is not None:
            try:
                return self.prs.slide_layouts[int(layout_index)]
            except Exception:
                pass

        layout_name = slide_config.get("template_layout")
        if layout_name:
            layout = self._find_layout([str(layout_name)])
            if layout is not None:
                return layout

        layout = self._find_layout(name_candidates)
        if layout is not None:
            return layout

        if self.theme_name == "template":
            auto = self._auto_select_layout(kind)
            if auto is not None:
                return auto

        return self._get_blank_layout()

    def _auto_select_layout(self, kind: str):
        def placeholder_types(layout) -> list:
            types = []
            for placeholder in layout.placeholders:
                try:
                    types.append(placeholder.placeholder_format.type)
                except Exception:
                    continue
            return types

        def has_title(types: list) -> bool:
            return PP_PLACEHOLDER.TITLE in types or PP_PLACEHOLDER.CENTER_TITLE in types

        def body_placeholders(layout) -> list:
            body_types = {
                PP_PLACEHOLDER.BODY,
                PP_PLACEHOLDER.OBJECT,
                PP_PLACEHOLDER.VERTICAL_BODY,
                PP_PLACEHOLDER.VERTICAL_OBJECT,
            }
            bodies = []
            for placeholder in layout.placeholders:
                try:
                    if placeholder.placeholder_format.type in body_types and hasattr(placeholder, "left"):
                        bodies.append(placeholder)
                except Exception:
                    continue
            return bodies

        def score(layout) -> int:
            types = placeholder_types(layout)
            title_ok = has_title(types)
            subtitle_ok = PP_PLACEHOLDER.SUBTITLE in types
            bodies = body_placeholders(layout)
            layout_name = str(getattr(layout, "name", "") or "").lower()

            points = 0
            if kind == "title":
                points += 10 if title_ok else 0
                points += 5 if subtitle_ok else 0
                points -= 3 * max(0, len(bodies) - 1)
                if "title slide" in layout_name:
                    points += 8
                if any(flag in layout_name for flag in ("chapter", "closing", "agenda", "thank you")):
                    points -= 16
            elif kind == "bullets":
                points += 8 if title_ok else 0
                points += 8 if len(bodies) >= 1 else 0
                points -= 2 * max(0, len(bodies) - 1)
                if "content slide" in layout_name:
                    points += 10
                if "agenda" in layout_name:
                    points -= 24
                if any(flag in layout_name for flag in ("closing", "thank you", "chapter")):
                    points -= 12
            elif kind in {"workflow", "kpi-cards", "infographic"}:
                # Prefer a title + a single large body area for diagrams/cards.
                points += 8 if title_ok else 0
                points += 6 if len(bodies) >= 1 else 0
                points -= 2 * max(0, len(bodies) - 1)
                if "content slide" in layout_name:
                    points += 8
                if any(flag in layout_name for flag in ("agenda", "closing", "chapter", "thank you")):
                    points -= 14
            elif kind == "two-column":
                points += 6 if title_ok else 0
                if len(bodies) >= 2:
                    slide_width = int(self.prs.slide_width)
                    midpoint = slide_width / 2
                    left_side = [b for b in bodies if int(b.left) < midpoint]
                    right_side = [b for b in bodies if int(b.left) >= midpoint]
                    points += 10 if left_side and right_side else 4
                else:
                    points -= 10
                if "agenda" in layout_name:
                    points -= 24
                if any(flag in layout_name for flag in ("closing", "chapter", "thank you")):
                    points -= 12
            elif kind == "image":
                points += 6 if title_ok else 0
                picture_ok = PP_PLACEHOLDER.PICTURE in types or PP_PLACEHOLDER.BITMAP in types
                points += 8 if picture_ok else 0
                if "content slide" in layout_name:
                    points += 4
                if any(flag in layout_name for flag in ("agenda", "closing", "chapter", "thank you")):
                    points -= 10
            elif kind == "chart":
                points += 6 if title_ok else 0
                chart_ok = PP_PLACEHOLDER.CHART in types
                points += 10 if chart_ok else 0
                points += 4 if len(bodies) >= 1 else 0
                if "content slide" in layout_name:
                    points += 6
                if any(flag in layout_name for flag in ("agenda", "closing", "chapter", "thank you")):
                    points -= 12
            elif kind == "blank":
                # Prefer truly blank or near-blank layouts
                points += 5 if not types else 0
                points -= len(types)
                if any(flag in layout_name for flag in ("closing", "chapter", "agenda", "thank you")):
                    points -= 8
            else:
                points += 0

            return points

        best_layout = None
        best_score = -10**9
        for layout in self.prs.slide_layouts:
            try:
                s = score(layout)
            except Exception:
                continue
            if s > best_score:
                best_layout = layout
                best_score = s

        return best_layout if best_score > 0 else None

    def _set_title_placeholder(self, slide, text: str) -> bool:
        title_shape = getattr(slide.shapes, "title", None)
        if not title_shape or not hasattr(title_shape, "text_frame"):
            return False
        title_shape.text_frame.text = text or ""
        title_shape.text_frame.word_wrap = True
        title_shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return True

    def _set_subtitle_placeholder(self, slide, text: str) -> bool:
        if not text:
            return False

        for placeholder in slide.placeholders:
            try:
                if placeholder.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE and hasattr(
                    placeholder, "text_frame"
                ):
                    placeholder.text_frame.text = text
                    placeholder.text_frame.word_wrap = True
                    placeholder.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    return True
            except Exception:
                continue

        return False

    def _get_best_body_placeholder(self, slide):
        candidates = self._get_body_placeholders(slide)
        if not candidates:
            return None

        return max(candidates, key=lambda shape: int(shape.width) * int(shape.height))

    def _get_body_placeholders(self, slide) -> list:
        body_types = {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.VERTICAL_BODY,
            PP_PLACEHOLDER.VERTICAL_OBJECT,
        }
        candidates = []
        for placeholder in slide.placeholders:
            try:
                if placeholder.placeholder_format.type in body_types:
                    candidates.append(placeholder)
            except Exception:
                continue
        return candidates

    def _fill_text_frame(self, text_frame, lines: list[str], as_bullets: bool) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        lines = lines or [""]

        # Template placeholders often use conservative defaults (small font, low vertical anchoring).
        # Tighten margins and anchor to the top so content uses the available space.
        try:
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            text_frame.margin_left = Inches(0.12)
            text_frame.margin_right = Inches(0.12)
            text_frame.margin_top = Inches(0.06)
            # Reserve space for template master elements (e.g., bottom-left logo).
            text_frame.margin_bottom = (
                Inches(0.32) if (self.template_path and self.theme_name == "template") else Inches(0.06)
            )
        except Exception:
            pass

        # Balanced text sizing to improve fit/readability in template placeholders.
        base_size = Pt(20 if as_bullets else 18)
        space_after = Pt(5 if as_bullets else 3)

        for i, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            cleaned = line.lstrip()
            if cleaned.startswith("•"):
                cleaned = cleaned.lstrip("•").lstrip()
            paragraph.text = cleaned
            if as_bullets:
                paragraph.level = 0
            else:
                paragraph.level = 0
            try:
                paragraph.font.size = base_size
                paragraph.space_after = space_after
            except Exception:
                pass

    def _fill_two_column_placeholders(self, slide, left_text: str, right_text: str) -> bool:
        left_text = self._normalize_newlines(left_text)
        right_text = self._normalize_newlines(right_text)

        placeholders = list(slide.placeholders)
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None:
            placeholders = [p for p in placeholders if p != title_shape]

        text_placeholders = [p for p in placeholders if hasattr(p, "text_frame")]
        if len(text_placeholders) < 2:
            return False

        slide_width = int(self.prs.slide_width)
        midpoint = slide_width / 2

        # Prefer largest text placeholders on opposite halves to avoid filling tiny/aux placeholders.
        def center_x(ph) -> float:
            return int(ph.left) + int(ph.width) / 2

        left_main = [p for p in text_placeholders if center_x(p) < midpoint]
        right_main = [p for p in text_placeholders if center_x(p) > midpoint]
        if left_main and right_main:
            left_primary = max(left_main, key=lambda p: int(p.width) * int(p.height))
            right_primary = max(right_main, key=lambda p: int(p.width) * int(p.height))
            self._fill_text_frame(
                left_primary.text_frame,
                left_text.split("\n") if left_text else [""],
                as_bullets=False,
            )
            self._fill_text_frame(
                right_primary.text_frame,
                right_text.split("\n") if right_text else [""],
                as_bullets=False,
            )
            return True

        left_side = [p for p in text_placeholders if int(p.left) < midpoint]
        right_side = [p for p in text_placeholders if int(p.left) >= midpoint]

        left_heading, left_body = self._pick_heading_and_body(left_side)
        right_heading, right_body = self._pick_heading_and_body(right_side)

        left_lines = left_text.split("\n") if left_text else [""]
        right_lines = right_text.split("\n") if right_text else [""]

        if left_heading and left_body and right_heading and right_body and len(text_placeholders) >= 4:
            left_heading.text_frame.text = left_lines[0] if left_lines else ""
            right_heading.text_frame.text = right_lines[0] if right_lines else ""

            self._fill_text_frame(left_body.text_frame, left_lines[1:] if len(left_lines) > 1 else [], as_bullets=True)
            self._fill_text_frame(
                right_body.text_frame, right_lines[1:] if len(right_lines) > 1 else [], as_bullets=True
            )
            return True

        if len(text_placeholders) >= 2:
            left_ph, right_ph = sorted(text_placeholders[:2], key=lambda p: int(p.left))
            self._fill_text_frame(left_ph.text_frame, left_lines, as_bullets=False)
            self._fill_text_frame(right_ph.text_frame, right_lines, as_bullets=False)
            return True

        return False

    def _pick_heading_and_body(self, side_placeholders: list) -> tuple[Optional[Any], Optional[Any]]:
        text_placeholders = [p for p in side_placeholders if hasattr(p, "text_frame")]
        if not text_placeholders:
            return None, None

        heading = min(text_placeholders, key=lambda p: int(p.top))
        body = max(text_placeholders, key=lambda p: int(p.height))
        return heading, body

    def _normalize_newlines(self, text: Optional[str]) -> str:
        if not text:
            return ""
        # Accept either "\n" newlines or literal "\\n" sequences in JSON.
        return text.replace("\r\n", "\n").replace("\\n", "\n")

    def _resolve_fs_path(self, path_like: Any) -> Optional[Path]:
        if not path_like:
            return None
        raw = str(path_like)

        # First try as-is (absolute or relative to cwd)
        p = Path(raw)
        if p.exists():
            return p

        # Try relative to the config directory when available
        if self.config_dir:
            p2 = Path(self.config_dir) / raw
            if p2.exists():
                return p2

        # WSL convenience: convert a Windows path (C:\...) to /mnt/c/...
        if re.match(r"^[A-Za-z]:\\\\", raw) or raw.startswith("\\\\"):
            try:
                result = subprocess.run(
                    ["wslpath", "-u", raw],
                    check=True,
                    capture_output=True,
                    text=True,
                )
                converted = result.stdout.strip()
                p3 = Path(converted)
                if p3.exists():
                    return p3
            except Exception:
                return p

        return p

    def _slugify(self, text: str, *, fallback: str = "image") -> str:
        cleaned = re.sub(r"[^A-Za-z0-9]+", "-", (text or "").strip()).strip("-").lower()
        return cleaned or fallback

    def _assets_output_dir(self) -> Path:
        if self.assets_dir:
            return self.assets_dir
        if self.config_dir:
            return Path(self.config_dir) / "_generated_assets"
        return Path.cwd() / "_generated_assets"

    def _color_to_pil_rgb(self, color: RGBColor | MSO_THEME_COLOR, *, default: Tuple[int, int, int]) -> Tuple[int, int, int]:
        if isinstance(color, MSO_THEME_COLOR):
            rgb = self._theme_to_rgb(color)
        else:
            rgb = color
        try:
            return tuple(int(c) for c in rgb)  # type: ignore[arg-type]
        except Exception:
            return default

    def _resolve_image_px_size(self, spec: dict) -> Tuple[int, int]:
        def as_int(value: Any) -> Optional[int]:
            try:
                return int(value)
            except Exception:
                return None

        width = as_int(spec.get("width_px") or spec.get("width") or spec.get("w"))
        height = as_int(spec.get("height_px") or spec.get("height") or spec.get("h"))
        if width and height:
            return max(64, width), max(64, height)

        aspect_raw = spec.get("aspect_ratio") or spec.get("aspect") or spec.get("size")
        aspect = None
        if aspect_raw:
            m = re.match(r"^\\s*(\\d+(?:\\.\\d+)?)\\s*:\\s*(\\d+(?:\\.\\d+)?)\\s*$", str(aspect_raw))
            if m:
                try:
                    a = float(m.group(1))
                    b = float(m.group(2))
                    if a > 0 and b > 0:
                        aspect = a / b
                except Exception:
                    aspect = None

        long_side = as_int(spec.get("long_side_px") or spec.get("max_px") or spec.get("px")) or 1600
        long_side = max(256, long_side)
        if not aspect:
            aspect = 16 / 9

        if aspect >= 1:
            w = long_side
            h = int(round(long_side / aspect))
        else:
            h = long_side
            w = int(round(long_side * aspect))

        return max(64, w), max(64, h)

    def _generate_image_asset(self, slide_config: dict) -> Optional[Path]:
        gen = slide_config.get("image_gen") or slide_config.get("generate_image") or slide_config.get("generate")
        if not isinstance(gen, dict):
            return None

        out_dir = self._assets_output_dir()
        out_dir.mkdir(parents=True, exist_ok=True)

        fmt = str(gen.get("format") or "png").strip().lower()
        if fmt == "jpeg":
            fmt = "jpg"
        if fmt not in {"png", "jpg"}:
            fmt = "png"

        image_id = str(gen.get("id") or self._slugify(str(slide_config.get("title") or "image"))).strip()
        image_id = self._slugify(image_id, fallback="image")
        out_path = out_dir / f"{image_id}.{fmt}"

        overwrite = bool(gen.get("overwrite", False))
        if out_path.exists() and not overwrite:
            return out_path

        # Palette for images: accept explicit RGB/hex or theme colors (resolved to RGB via template theme).
        palette_cfg = getattr(self, "presentation_config", {}) or {}
        pres_img_section = self._coerce_palette_section(palette_cfg.get("palette"), "image")
        slide_img_section = self._coerce_palette_section(slide_config.get("palette"), "image")
        img_palette = {**pres_img_section, **slide_img_section}

        bg_color = self._parse_color(img_palette.get("bg") or gen.get("bg") or gen.get("background") or "#FFFFFF")
        accent_color = self._parse_color(img_palette.get("accent") or gen.get("accent") or "ACCENT_1")
        text_color = self._parse_color(img_palette.get("text") or gen.get("text_color") or "TEXT_1")

        if bg_color is None:
            bg_color = RGBColor(255, 255, 255)
        if accent_color is None:
            accent_color = MSO_THEME_COLOR.ACCENT_1 if self.theme_name == "template" else RGBColor(2, 132, 199)
        if text_color is None:
            text_color = MSO_THEME_COLOR.TEXT_1 if self.theme_name == "template" else RGBColor(15, 23, 42)

        try:
            self._render_generated_image(
                gen,
                out_path=out_path,
                bg=bg_color,
                accent=accent_color,
                text=text_color,
            )
            return out_path
        except Exception as e:
            print(f"⚠️  Image generation failed for {out_path}: {e}", file=sys.stderr)
            return None

    def _render_generated_image(
        self,
        spec: dict,
        *,
        out_path: Path,
        bg: RGBColor | MSO_THEME_COLOR,
        accent: RGBColor | MSO_THEME_COLOR,
        text: RGBColor | MSO_THEME_COLOR,
    ) -> None:
        try:
            from PIL import Image, ImageDraw, ImageFilter, ImageFont
        except Exception as e:
            raise RuntimeError("Pillow is required for image generation. Install scripts/requirements.txt") from e

        w_px, h_px = self._resolve_image_px_size(spec)
        kind = str(spec.get("kind") or spec.get("style") or "infographic").strip().lower()
        prompt = str(spec.get("prompt") or spec.get("text") or "").strip()
        seed_raw = spec.get("seed")
        if seed_raw is None:
            seed_raw = zlib.crc32(prompt.encode("utf-8")) if prompt else 0
        try:
            seed = int(seed_raw)
        except Exception:
            seed = 0
        rng = random.Random(seed)

        bg_rgb = self._color_to_pil_rgb(bg, default=(255, 255, 255))
        accent_rgb = self._color_to_pil_rgb(accent, default=(2, 132, 199))
        text_rgb = self._color_to_pil_rgb(text, default=(15, 23, 42))

        img = Image.new("RGB", (w_px, h_px), bg_rgb)
        draw = ImageDraw.Draw(img)

        def try_font(size: int) -> ImageFont.ImageFont:
            candidates = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            ]
            for path in candidates:
                try:
                    if Path(path).exists():
                        return ImageFont.truetype(path, size=size)
                except Exception:
                    continue
            return ImageFont.load_default()

        title_font = try_font(max(14, int(min(w_px, h_px) * 0.05)))
        body_font = try_font(max(12, int(min(w_px, h_px) * 0.03)))

        if kind in {"photo", "realistic", "image"}:
            # Abstract "photo-like" background (gradient + noise) as a placeholder.
            start = bg_rgb
            end = tuple(max(0, min(c + rng.randint(-30, 30), 255)) for c in accent_rgb)
            for y in range(h_px):
                t = y / max(1, h_px - 1)
                row = tuple(int(start[i] * (1 - t) + end[i] * t) for i in range(3))
                draw.line([(0, y), (w_px, y)], fill=row)
            for _ in range(int(w_px * h_px * 0.003)):
                x = rng.randint(0, w_px - 1)
                y = rng.randint(0, h_px - 1)
                jitter = rng.randint(-18, 18)
                px = tuple(max(0, min(int(c) + jitter, 255)) for c in img.getpixel((x, y)))
                img.putpixel((x, y), px)
            img = img.filter(ImageFilter.GaussianBlur(radius=1.2))

        elif kind in {"sketch", "handdrawn"}:
            # Simple sketch effect: jittered lines + light paper grain.
            for _ in range(12):
                x1, y1 = rng.randint(0, w_px), rng.randint(0, h_px)
                x2, y2 = rng.randint(0, w_px), rng.randint(0, h_px)
                for offset in (-2, -1, 0, 1, 2):
                    draw.line(
                        [(x1 + rng.randint(-3, 3), y1 + rng.randint(-3, 3)), (x2 + rng.randint(-3, 3), y2 + rng.randint(-3, 3))],
                        fill=tuple(max(0, min(c + offset, 255)) for c in text_rgb),
                        width=2,
                    )
            for _ in range(int(w_px * h_px * 0.002)):
                x = rng.randint(0, w_px - 1)
                y = rng.randint(0, h_px - 1)
                px = img.getpixel((x, y))
                img.putpixel((x, y), tuple(max(0, min(c + rng.randint(-8, 8), 255)) for c in px))

        elif kind in {"workflow", "infographic-workflow"}:
            steps = spec.get("steps") or []
            if not isinstance(steps, list):
                steps = []
            steps = [str(s.get("title") if isinstance(s, dict) else s).strip() for s in steps if str(s).strip()]
            if not steps:
                steps = ["Step 1", "Step 2", "Step 3"]
            pad = int(min(w_px, h_px) * 0.06)
            box_h = int(h_px * 0.22)
            box_w = int((w_px - 2 * pad) / max(1, len(steps))) - int(pad * 0.15)
            y0 = int(h_px * 0.38)
            for i, label in enumerate(steps):
                x0 = pad + i * (box_w + int(pad * 0.15))
                rect = [x0, y0, x0 + box_w, y0 + box_h]
                draw.rounded_rectangle(rect, radius=18, fill=tuple(int(c * 0.96) for c in bg_rgb), outline=accent_rgb, width=4)
                bbox = draw.textbbox((0, 0), label, font=body_font)
                tw = bbox[2] - bbox[0]
                th = bbox[3] - bbox[1]
                draw.text((x0 + (box_w - tw) / 2, y0 + (box_h - th) / 2), label, fill=text_rgb, font=body_font)
                if i < len(steps) - 1:
                    ax1 = x0 + box_w + int(pad * 0.03)
                    ax2 = ax1 + int(pad * 0.35)
                    ay = y0 + box_h // 2
                    draw.line([(ax1, ay), (ax2, ay)], fill=accent_rgb, width=5)
                    draw.polygon([(ax2, ay), (ax2 - 12, ay - 9), (ax2 - 12, ay + 9)], fill=accent_rgb)

        elif kind in {"ui-dashboard", "react-ui", "dashboard-ui", "app-screenshot"}:
            def mix(c1: tuple[int, int, int], c2: tuple[int, int, int], t: float) -> tuple[int, int, int]:
                t = max(0.0, min(1.0, t))
                return tuple(int(c1[i] * (1 - t) + c2[i] * t) for i in range(3))

            def luma(rgb: tuple[int, int, int]) -> float:
                r, g, b = rgb
                return 0.299 * r + 0.587 * g + 0.114 * b

            bg_is_dark = luma(bg_rgb) < 120
            if bg_is_dark:
                top_bg = mix(bg_rgb, (255, 255, 255), 0.03)
                frame_bg = mix(bg_rgb, (255, 255, 255), 0.05)
                surface = mix(bg_rgb, (255, 255, 255), 0.10)
                surface_alt = mix(bg_rgb, (255, 255, 255), 0.15)
                text_primary = mix(text_rgb, (255, 255, 255), 0.05)
                text_secondary = mix(text_rgb, (200, 215, 235), 0.35)
            else:
                top_bg = mix(bg_rgb, (0, 0, 0), 0.03)
                frame_bg = mix(bg_rgb, (0, 0, 0), 0.05)
                surface = mix(bg_rgb, (0, 0, 0), 0.06)
                surface_alt = mix(bg_rgb, (0, 0, 0), 0.10)
                text_primary = mix(text_rgb, (10, 20, 35), 0.18)
                text_secondary = mix(text_rgb, (60, 80, 105), 0.30)

            for y in range(h_px):
                t = y / max(1, h_px - 1)
                row = mix(top_bg, frame_bg, t)
                draw.line([(0, y), (w_px, y)], fill=row)

            pad = int(min(w_px, h_px) * 0.035)
            frame = [pad, pad, w_px - pad, h_px - pad]
            draw.rounded_rectangle(frame, radius=max(16, int(pad * 0.7)), fill=frame_bg, outline=accent_rgb, width=3)

            fx0, fy0, fx1, fy1 = frame
            fw = fx1 - fx0
            fh = fy1 - fy0
            gap = max(10, int(min(w_px, h_px) * 0.012))
            top_h = int(fh * 0.12)
            side_w = int(fw * 0.21)

            top_rect = [fx0 + 2, fy0 + 2, fx1 - 2, fy0 + top_h]
            draw.rounded_rectangle(top_rect, radius=max(10, int(pad * 0.35)), fill=surface_alt)

            side_rect = [fx0 + 2, fy0 + top_h + gap, fx0 + side_w, fy1 - 2]
            draw.rounded_rectangle(side_rect, radius=max(10, int(pad * 0.35)), fill=surface)

            brand = str(spec.get("app_title") or "AutoBidding Console").strip()
            pfont = try_font(max(18, int(min(w_px, h_px) * 0.028)))
            sfont = try_font(max(12, int(min(w_px, h_px) * 0.020)))
            tiny_font = try_font(max(10, int(min(w_px, h_px) * 0.016)))

            draw.ellipse(
                [fx0 + gap, fy0 + gap // 2, fx0 + gap + int(top_h * 0.52), fy0 + gap // 2 + int(top_h * 0.52)],
                fill=accent_rgb,
            )
            draw.text((fx0 + gap + int(top_h * 0.7), fy0 + gap // 2), brand, fill=text_primary, font=sfont)

            search_w = int(fw * 0.34)
            search_h = int(top_h * 0.52)
            search_x = fx0 + int(fw * 0.34)
            search_y = fy0 + (top_h - search_h) // 2
            draw.rounded_rectangle(
                [search_x, search_y, search_x + search_w, search_y + search_h],
                radius=max(8, int(search_h * 0.45)),
                fill=surface,
                outline=mix(accent_rgb, frame_bg, 0.65),
                width=2,
            )
            draw.text((search_x + int(search_h * 0.38), search_y + int(search_h * 0.20)), "Search tenders...", fill=text_secondary, font=tiny_font)

            chip_w = int(top_h * 0.56)
            chip_h = chip_w
            chip_x = fx1 - gap - chip_w
            chip_y = fy0 + (top_h - chip_h) // 2
            draw.rounded_rectangle(
                [chip_x, chip_y, chip_x + chip_w, chip_y + chip_h],
                radius=max(8, int(chip_w * 0.35)),
                fill=mix(accent_rgb, frame_bg, 0.78),
                outline=accent_rgb,
                width=2,
            )

            nav_items = spec.get("sidebar_items") or [
                "Workspace",
                "Opportunities",
                "Research",
                "Drafting",
                "Review",
                "Exports",
            ]
            if not isinstance(nav_items, list):
                nav_items = ["Workspace", "Opportunities", "Research", "Drafting", "Review", "Exports"]
            nav_items = [str(x).strip() for x in nav_items if str(x).strip()]
            if not nav_items:
                nav_items = ["Workspace", "Opportunities", "Research", "Drafting", "Review", "Exports"]
            active_idx = int(spec.get("active_nav_index") or 1)
            nav_top = side_rect[1] + gap
            nav_line_h = max(20, int((side_rect[3] - nav_top - gap) / max(len(nav_items), 6)))
            for i, label in enumerate(nav_items[:8]):
                y = nav_top + i * nav_line_h
                is_active = i == active_idx
                if is_active:
                    draw.rounded_rectangle(
                        [side_rect[0] + gap // 2, y - 2, side_rect[2] - gap // 2, y + nav_line_h - 4],
                        radius=max(8, int(nav_line_h * 0.35)),
                        fill=mix(accent_rgb, frame_bg, 0.75),
                    )
                ix = side_rect[0] + gap
                iy = y + int(nav_line_h * 0.24)
                isize = int(nav_line_h * 0.36)
                draw.rounded_rectangle(
                    [ix, iy, ix + isize, iy + isize],
                    radius=max(3, int(isize * 0.24)),
                    fill=mix(accent_rgb, frame_bg, 0.45 if is_active else 0.65),
                )
                draw.text((ix + isize + gap // 2, y + int(nav_line_h * 0.16)), label, fill=text_primary, font=tiny_font)

            content_x0 = side_rect[2] + gap
            content_y0 = side_rect[1]
            content_x1 = fx1 - gap
            content_y1 = fy1 - gap
            content_w = content_x1 - content_x0
            content_h = content_y1 - content_y0

            kpi_h = int(content_h * 0.23)
            kpi_gap = gap
            kpi_y = content_y0
            kpi_cards = spec.get("kpi_cards") or [
                {"title": "Win Rate", "value": "41%", "delta": "+5.2%"},
                {"title": "Avg Cycle", "value": "2.8d", "delta": "-18%"},
                {"title": "Active Bids", "value": "27", "delta": "+4"},
            ]
            if not isinstance(kpi_cards, list):
                kpi_cards = []
            if len(kpi_cards) < 3:
                kpi_cards = [
                    {"title": "Win Rate", "value": "41%", "delta": "+5.2%"},
                    {"title": "Avg Cycle", "value": "2.8d", "delta": "-18%"},
                    {"title": "Active Bids", "value": "27", "delta": "+4"},
                ]

            kpi_cols = 3
            card_w = int((content_w - (kpi_cols - 1) * kpi_gap) / kpi_cols)
            for i in range(kpi_cols):
                card = kpi_cards[i] if i < len(kpi_cards) else {}
                cx0 = content_x0 + i * (card_w + kpi_gap)
                cy0 = kpi_y
                cx1 = cx0 + card_w
                cy1 = cy0 + kpi_h
                draw.rounded_rectangle([cx0, cy0, cx1, cy1], radius=max(10, int(kpi_h * 0.18)), fill=surface_alt, outline=mix(accent_rgb, frame_bg, 0.55), width=2)
                draw.text((cx0 + gap, cy0 + int(kpi_h * 0.14)), str(card.get("title") or f"Metric {i+1}"), fill=text_secondary, font=tiny_font)
                draw.text((cx0 + gap, cy0 + int(kpi_h * 0.42)), str(card.get("value") or "--"), fill=text_primary, font=pfont)
                draw.text((cx0 + gap, cy0 + int(kpi_h * 0.72)), str(card.get("delta") or ""), fill=accent_rgb, font=tiny_font)

            lower_y0 = kpi_y + kpi_h + gap
            lower_h = content_y1 - lower_y0
            left_w = int(content_w * 0.57)
            right_w = content_w - left_w - gap

            board = [content_x0, lower_y0, content_x0 + left_w, lower_y0 + lower_h]
            preview = [board[2] + gap, lower_y0, board[2] + gap + right_w, lower_y0 + lower_h]
            draw.rounded_rectangle(board, radius=max(10, int(lower_h * 0.08)), fill=surface, outline=mix(accent_rgb, frame_bg, 0.62), width=2)
            draw.rounded_rectangle(preview, radius=max(10, int(lower_h * 0.08)), fill=surface, outline=mix(accent_rgb, frame_bg, 0.62), width=2)

            draw.text((board[0] + gap, board[1] + gap), "Workflow Board", fill=text_primary, font=sfont)
            lanes = spec.get("lanes") or ["Research", "Drafting", "Review"]
            if not isinstance(lanes, list):
                lanes = ["Research", "Drafting", "Review"]
            lanes = [str(x).strip() for x in lanes if str(x).strip()]
            if len(lanes) < 3:
                lanes = ["Research", "Drafting", "Review"]
            lane_gap = gap
            lane_top = board[1] + int(lower_h * 0.16)
            lane_h = board[3] - lane_top - gap
            lane_w = int((left_w - 2 * lane_gap - 2 * gap) / 3)
            for i in range(3):
                lx0 = board[0] + gap + i * (lane_w + lane_gap)
                lx1 = lx0 + lane_w
                draw.rounded_rectangle([lx0, lane_top, lx1, lane_top + lane_h], radius=max(8, int(lane_h * 0.06)), fill=mix(surface_alt, frame_bg, 0.35), outline=mix(accent_rgb, frame_bg, 0.72), width=2)
                lane_label = lanes[i] if i < len(lanes) else f"Lane {i+1}"
                draw.text((lx0 + int(lane_w * 0.08), lane_top + int(lane_h * 0.03)), lane_label, fill=text_secondary, font=tiny_font)
                card_h = int(lane_h * 0.22)
                for j in range(2):
                    cy0 = lane_top + int(lane_h * 0.17) + j * (card_h + int(lane_h * 0.08))
                    cy1 = cy0 + card_h
                    draw.rounded_rectangle(
                        [lx0 + int(lane_w * 0.08), cy0, lx1 - int(lane_w * 0.08), cy1],
                        radius=max(6, int(card_h * 0.25)),
                        fill=surface_alt,
                        outline=mix(accent_rgb, frame_bg, 0.58),
                        width=2,
                    )
                    draw.line(
                        [(lx0 + int(lane_w * 0.12), cy0 + int(card_h * 0.35)), (lx1 - int(lane_w * 0.12), cy0 + int(card_h * 0.35))],
                        fill=text_secondary,
                        width=2,
                    )
                    draw.line(
                        [(lx0 + int(lane_w * 0.12), cy0 + int(card_h * 0.62)), (lx0 + int(lane_w * 0.58), cy0 + int(card_h * 0.62))],
                        fill=text_secondary,
                        width=2,
                    )

            draw.text((preview[0] + gap, preview[1] + gap), "Proposal Preview", fill=text_primary, font=sfont)
            table_top = preview[1] + int(lower_h * 0.16)
            table = [preview[0] + gap, table_top, preview[2] - gap, preview[3] - gap]
            draw.rounded_rectangle(table, radius=max(8, int(lower_h * 0.05)), fill=surface_alt, outline=mix(accent_rgb, frame_bg, 0.66), width=2)

            headers = spec.get("table_headers") or ["Section", "Owner", "Status"]
            if not isinstance(headers, list):
                headers = ["Section", "Owner", "Status"]
            headers = [str(x).strip() for x in headers if str(x).strip()]
            if len(headers) < 2:
                headers = ["Section", "Owner", "Status"]
            rows = spec.get("table_rows") or [
                ["Executive Summary", "Analyst A", "Ready"],
                ["Pricing Approach", "Analyst B", "In Review"],
                ["Delivery Model", "Analyst C", "Draft"],
                ["Compliance Matrix", "Analyst D", "Ready"],
            ]
            if not isinstance(rows, list):
                rows = []
            rows_clean = []
            for row in rows:
                if isinstance(row, (list, tuple)):
                    rows_clean.append([str(x).strip() for x in row])
            if not rows_clean:
                rows_clean = [
                    ["Executive Summary", "Analyst A", "Ready"],
                    ["Pricing Approach", "Analyst B", "In Review"],
                    ["Delivery Model", "Analyst C", "Draft"],
                    ["Compliance Matrix", "Analyst D", "Ready"],
                ]

            cols = max(2, len(headers))
            tx0, ty0, tx1, ty1 = table
            th = max(18, int((ty1 - ty0) * 0.16))
            draw.rectangle([tx0 + 1, ty0 + 1, tx1 - 1, ty0 + th], fill=mix(surface, frame_bg, 0.20))

            col_w = (tx1 - tx0) / cols
            for i in range(1, cols):
                x = int(tx0 + i * col_w)
                draw.line([(x, ty0), (x, ty1)], fill=mix(accent_rgb, frame_bg, 0.70), width=1)
            draw.line([(tx0, ty0 + th), (tx1, ty0 + th)], fill=mix(accent_rgb, frame_bg, 0.60), width=1)

            for i, header in enumerate(headers[:cols]):
                hx = int(tx0 + i * col_w + gap * 0.6)
                draw.text((hx, ty0 + int(th * 0.22)), header, fill=text_secondary, font=tiny_font)

            body_rows = min(5, len(rows_clean))
            if body_rows > 0:
                row_h = max(14, int((ty1 - (ty0 + th)) / body_rows))
                for r in range(body_rows):
                    y = int(ty0 + th + r * row_h)
                    if r % 2 == 1:
                        draw.rectangle([tx0 + 1, y, tx1 - 1, y + row_h], fill=mix(surface_alt, frame_bg, 0.24))
                    draw.line([(tx0, y + row_h), (tx1, y + row_h)], fill=mix(accent_rgb, frame_bg, 0.74), width=1)
                    row_data = rows_clean[r]
                    for c in range(cols):
                        cell = row_data[c] if c < len(row_data) else ""
                        cx = int(tx0 + c * col_w + gap * 0.6)
                        draw.text((cx, y + int(row_h * 0.20)), cell, fill=text_primary, font=tiny_font)

        else:
            # Generic infographic: headline + a few accent blocks.
            blocks = int(spec.get("blocks") or 4)
            blocks = max(2, min(blocks, 6))
            pad = int(min(w_px, h_px) * 0.06)
            block_h = int(h_px * 0.18)
            for i in range(blocks):
                x0 = pad
                y0 = int(h_px * 0.35) + i * (block_h + int(pad * 0.25))
                x1 = w_px - pad
                y1 = y0 + block_h
                if y1 > h_px - pad:
                    break
                tint = tuple(max(0, min(c + rng.randint(-20, 20), 255)) for c in accent_rgb)
                draw.rounded_rectangle([x0, y0, x1, y1], radius=16, outline=tint, width=4)

        if prompt:
            bbox = draw.textbbox((0, 0), prompt, font=title_font)
            tw = bbox[2] - bbox[0]
            th = bbox[3] - bbox[1]
            draw.text(((w_px - tw) / 2, int(h_px * 0.10) - th / 2), prompt, fill=text_rgb, font=title_font)

        out_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(str(out_path))

    def _add_picture_contain(self, slide, *, img_path: Path, x, y, cx, cy) -> None:
        left, top, w, h = self._compute_contain_geometry(img_path=img_path, x=x, y=y, cx=cx, cy=cy)
        try:
            slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
        except Exception:
            slide.shapes.add_picture(str(img_path), x, y, width=cx)

    def _compute_contain_geometry(self, *, img_path: Path, x, y, cx, cy) -> tuple[int, int, int, int]:
        try:
            from PIL import Image
        except Exception:
            return int(x), int(y), int(cx), int(cy)

        try:
            with Image.open(img_path) as im:
                iw, ih = im.size
        except Exception:
            return int(x), int(y), int(cx), int(cy)

        try:
            box_w = float(int(cx))
            box_h = float(int(cy))
            if box_w <= 0 or box_h <= 0:
                return int(x), int(y), int(cx), int(cy)
            ratio = iw / max(1, ih)
            box_ratio = box_w / box_h
            if ratio >= box_ratio:
                w = box_w
                h = box_w / ratio
            else:
                h = box_h
                w = box_h * ratio
            left = int(int(x) + (box_w - w) / 2)
            top = int(int(y) + (box_h - h) / 2)
            return left, top, int(w), int(h)
        except Exception:
            return int(x), int(y), int(cx), int(cy)

    def _maybe_autocrop_image(self, img_path: Path, config: dict) -> Path:
        if not bool(config.get("image_autocrop")):
            return img_path
        if not self.assets_dir:
            return img_path

        try:
            from PIL import Image, ImageChops
        except Exception:
            return img_path

        out_path = self.assets_dir / f"{img_path.stem}-autocrop.png"
        try:
            out_path.parent.mkdir(parents=True, exist_ok=True)
        except Exception:
            return img_path

        try:
            with Image.open(img_path) as im_raw:
                im = im_raw.convert("RGB")
        except Exception:
            return img_path

        try:
            bg = Image.new("RGB", im.size, (255, 255, 255))
            diff = ImageChops.difference(im, bg).convert("L")
            mask = diff.point(lambda p: 255 if p > 12 else 0)
            bbox = mask.getbbox()
            if not bbox:
                return img_path
            l, t, r, b = bbox
            pad = max(10, int(min(im.size) * 0.02))
            l = max(0, l - pad)
            t = max(0, t - pad)
            r = min(im.size[0], r + pad)
            b = min(im.size[1], b + pad)
            cropped = im.crop((l, t, r, b))
            cropped.save(out_path)
            return out_path
        except Exception:
            return img_path

    def _add_chart_slide(self, config: dict) -> None:
        layout = self._resolve_layout(
            config,
            kind="chart",
            name_candidates=["chart", "title and content", "title and text", "content"],
        )
        slide = self.prs.slides.add_slide(layout)
        self._set_slide_background(slide)

        title = config.get("title", "")
        if title:
            if self.template_path and self.theme_name == "template":
                if not self._set_title_placeholder(slide, title):
                    self._add_title_textbox(slide, {"title": title}, y_pos=0.25)
            else:
                self._add_title_textbox(slide, {"title": title})

        chart_type_name = str(config.get("chart_type") or config.get("type") or "column").strip().lower()
        chart_type = self.CHART_TYPES.get(chart_type_name, XL_CHART_TYPE.COLUMN_CLUSTERED)

        categories = config.get("categories") or config.get("labels") or []
        series = config.get("series") or []

        if not categories and isinstance(series, list) and series:
            first = series[0] if isinstance(series[0], dict) else None
            values = (first or {}).get("values") if first else None
            if isinstance(values, list):
                categories = [f"Item {i+1}" for i in range(len(values))]

        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]

        has_series = False
        if isinstance(series, list) and series:
            for s in series:
                if not isinstance(s, dict):
                    continue
                name = str(s.get("name") or "")
                values = s.get("values") or []
                numeric_values = []
                for v in values:
                    try:
                        numeric_values.append(float(v))
                    except Exception:
                        numeric_values.append(0.0)
                chart_data.add_series(name, numeric_values)
                has_series = True

        if not has_series:
            # Fallback: empty chart with a single series.
            chart_data.add_series("Series", [0.0 for _ in chart_data.categories])

        chart = None
        if self.template_path and self.theme_name == "template":
            chart_ph = next((ph for ph in slide.placeholders if hasattr(ph, "insert_chart")), None)
            if chart_ph:
                graphic_frame = chart_ph.insert_chart(chart_type, chart_data)
                chart = graphic_frame.chart

        if chart is None:
            x, y, cx, cy = self._get_content_box(slide, has_title=bool(title))
            chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart

        legend = config.get("legend", True)
        chart.has_legend = bool(legend)
        if chart.has_legend:
            legend_position = str(config.get("legend_position") or "right").strip().lower()
            if legend_position in self.LEGEND_POSITIONS:
                chart.legend.position = self.LEGEND_POSITIONS[legend_position]
            chart.legend.include_in_layout = False

        style = config.get("style")
        if style is not None:
            try:
                chart.chart_style = int(style)
            except Exception:
                pass



from pptxgen.cli import run_cli
from pptxgen.qa_pipeline import export_snapshots as _export_snapshots
from pptxgen.qa_pipeline import run_qa_pipeline as _run_qa_pipeline


def main() -> None:
    run_cli(PPTXGenerator)


if __name__ == "__main__":
    main()
